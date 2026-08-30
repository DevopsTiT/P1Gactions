#!/usr/bin/env python3
"""Combine seq 13 (narrative) + seq 15 (design diagrams) + seq 17 (Teams POC)
into one whole PPT. Clean footers — no AXA / ARB Template branding."""

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x2C, 0x5C)
LIGHT_BLUE = RGBColor(0xD6, 0xE6, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x55, 0x55)
BLACK = RGBColor(0x22, 0x22, 0x22)
GREEN = RGBColor(0x1E, 0x84, 0x4E)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
LIGHT_ORANGE = RGBColor(0xFD, 0xE8, 0xD0)
NAVY_MID = RGBColor(0x14, 0x3D, 0x7A)

FOOTER = "Dynatrace to PagerDuty Four AI Agents — Whole Pack"
SRC_ROOT = Path("/Users/k/Learnings/AIProject/CursorFiles/Daily Files/2026-08-31")
SRC_13 = SRC_ROOT / "13-dynatrace-pagerduty-four-agents-ppt/13-Dynatrace-PagerDuty-Four-AI-Agents-ARB-style.pptx"
SRC_15 = SRC_ROOT / "15-dynatrace-pd-four-agents-design-diagrams/15-Dynatrace-PagerDuty-Four-Agents-Design-Diagrams.pptx"
SRC_17 = SRC_ROOT / "17-pagerduty-four-agents-teams-only-poc-ppt/17-PagerDuty-Four-Agents-Teams-Only-POC.pptx"

# Skip source title slides when a section divider already introduces the part
SKIP_13 = {0}  # title
SKIP_15 = {0}  # title
SKIP_17 = {0}  # title


def set_run(run, text, size=12, bold=False, color=BLACK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_textbox(slide, left, top, width, height, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color)
    return box


def add_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def copy_slide(dest_prs, source_slide):
    """Deep-copy a slide's shape tree into dest (works for shape/table/text decks)."""
    dest_slide = dest_prs.slides.add_slide(dest_prs.slide_layouts[6])
    # Remove default empty shapes from blank layout
    sp_tree = dest_slide.shapes._spTree
    for child in list(sp_tree):
        tag = child.tag
        if tag.endswith("}sp") or tag.endswith("}pic") or tag.endswith("}graphicFrame") or tag.endswith("}cxnSp") or tag.endswith("}grpSp"):
            sp_tree.remove(child)

    for shape in source_slide.shapes:
        el = deepcopy(shape.element)
        sp_tree.insert_element_before(el, "p:extLst")

    # Copy notes if present
    if source_slide.has_notes_slide:
        notes_text = source_slide.notes_slide.notes_text_frame.text
        if notes_text.strip():
            dest_slide.notes_slide.notes_text_frame.text = notes_text
    return dest_slide


def scrub_and_footer(slide, page, total):
    """Remove old footers / AXA-ARB phrases; stamp clean footer + page."""
    bad_phrases = (
        "AXA Japan ARB Template style",
        "ARB-style",
        "ARB style",
        "OneAXA",
        "Japan ARB",
        "LAST UPDATED on",
        "LAST UPDATED 2026",
    )
    old_footer_markers = (
        "AXA Japan ARB Template style",
        "ARB-style ·",
        "PagerDuty Four Agents — Teams-Only POC",
        "Dynatrace to PagerDuty Four AI Agents — Design",
        "Dynatrace → PagerDuty Four AI Agents · Design",
        "ARB-style working deck",
    )
    for shape in list(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = (shape.text_frame.text or "").strip()
        try:
            top = shape.top or 0
        except Exception:
            top = 0
        near_bottom = top > Inches(6.85)
        looks_like_page_num = near_bottom and text.isdigit()
        looks_like_old_footer = near_bottom and (
            any(m in text for m in old_footer_markers)
            or (len(text) < 90 and (" / " in text or text.replace(" ", "").isdigit()))
        )
        if looks_like_page_num or looks_like_old_footer:
            shape._element.getparent().remove(shape._element)
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text:
                    continue
                new = run.text
                for b in bad_phrases:
                    new = new.replace(b, "")
                new = new.replace(" ·  · ", " · ").strip(" ·")
                if new != run.text:
                    run.text = new

    add_textbox(slide, Inches(0.3), Inches(7.15), Inches(10), Inches(0.25), FOOTER, size=9, color=GRAY)
    add_textbox(
        slide,
        Inches(11.2),
        Inches(7.15),
        Inches(1.8),
        Inches(0.25),
        f"{page} / {total}",
        size=10,
        bold=True,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
    )


def section_divider(prs, part, title, bullets, fill_accent=TEAL):
    s = blank(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_textbox(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.4), part, size=16, bold=True, color=fill_accent)
    add_textbox(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.8), title, size=32, bold=True, color=WHITE)
    y = Inches(3.5)
    for b in bullets:
        add_textbox(s, Inches(0.7), y, Inches(11.5), Inches(0.4), f"• {b}", size=16, color=LIGHT_BLUE)
        y += Inches(0.45)
    return s


def master_title(prs):
    s = blank(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.4), "Combined architecture + design + Teams POC pack", size=14, bold=True, color=LIGHT_BLUE)
    add_textbox(s, Inches(0.5), Inches(1.7), Inches(12), Inches(0.7), "Dynatrace to PagerDuty", size=36, bold=True, color=WHITE)
    add_textbox(s, Inches(0.5), Inches(2.5), Inches(12), Inches(0.5), "Four AI Agents — Whole Deck", size=26, bold=True, color=LIGHT_BLUE)
    add_textbox(
        s,
        Inches(0.5),
        Inches(3.2),
        Inches(12),
        Inches(0.4),
        "SRE  |  Scribe  |  Shift  |  Insights",
        size=16,
        color=WHITE,
    )

    card_fill = RGBColor(0x12, 0x3A, 0x6B)
    for i, (title, body) in enumerate(
        [
            ("Part A — Narrative", "Architecture summary\nAS-IS / TO-BE story\nAgent detail + cost/safety"),
            ("Part B — Design diagrams", "Solution context boxes\nImpact / Alignment / HLD\nAgent hub + checklist"),
            ("Part C — Teams-only POC", "§0 enablement\nSRE / Scribe / Shift / Insights\nClick-path runbook"),
        ]
    ):
        left = Inches(0.5 + i * 4.2)
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(4.2), Inches(4.0), Inches(2.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = card_fill
        shape.line.color.rgb = LIGHT_BLUE
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_top = Inches(0.12)
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        set_run(r0, title, size=14, bold=True, color=WHITE)
        for line in body.split("\n"):
            p = tf.add_paragraph()
            r = p.add_run()
            set_run(r, line, size=12, color=LIGHT_BLUE)
    return s


def agenda_slide(prs, counts):
    s = blank(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.75), NAVY)
    add_textbox(s, Inches(0.35), Inches(0.18), Inches(12.5), Inches(0.45), "Whole Pack — Agenda", size=20, bold=True, color=WHITE)
    rows = [
        ("Part", "Source", "What you get", "Slides (approx)"),
        ("A", "seq 13 narrative PPT", "Summary sheet, AS-IS/TO-BE story, agents, enablement", str(counts["a"])),
        ("B", "seq 15 design diagrams PPT", "Visual Solution Context, Impact, Alignment, HLD, hub", str(counts["b"])),
        ("C", "seq 17 Teams-only POC PPT", "Full §0–§4 click path for Teams-only orgs", str(counts["c"])),
    ]
    from pptx.enum.text import MSO_ANCHOR
    from pptx.util import Inches as I

    table_shape = s.shapes.add_table(len(rows), 4, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.45 * len(rows)))
    table = table_shape.table
    widths = [Inches(1.2), Inches(3.8), Inches(5.5), Inches(2.0)]
    for i, w in enumerate(widths):
        table.columns[i].width = w
    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            is_h = r_idx == 0
            set_run(run, cell_text, size=12, bold=is_h, color=WHITE if is_h else BLACK)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if is_h else (RGBColor(0xF2, 0xF4, 0xF7) if r_idx % 2 else WHITE)

    add_textbox(
        s,
        Inches(0.4),
        Inches(4.0),
        Inches(12.5),
        Inches(2.2),
        "How to use\n"
        "• Design review / board discussion → Part A then Part B\n"
        "• Hands-on Teams-only POC afternoon → Part C\n"
        "• Safe demo rule everywhere: test Service + Level-1 test schedule — never prod pages\n"
        "• Footer on every slide is clean (no board-template branding)",
        size=14,
        color=BLACK,
    )
    return s


def append_from(src_path, dest_prs, skip_indices):
    src = Presentation(str(src_path))
    copied = 0
    for i, slide in enumerate(src.slides):
        if i in skip_indices:
            continue
        copy_slide(dest_prs, slide)
        copied += 1
    return copied


def build(out: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    master_title(prs)
    # temporary agenda; will rebuild counts after copy — insert placeholder then we fix by building agenda after
    # Actually build parts first into a list of slide groups by counting

    # Count planned content slides
    n13 = len(Presentation(str(SRC_13)).slides) - len(SKIP_13)
    n15 = len(Presentation(str(SRC_15)).slides) - len(SKIP_15)
    n17 = len(Presentation(str(SRC_17)).slides) - len(SKIP_17)
    # + master title + agenda + 3 section dividers
    agenda_slide(prs, {"a": n13, "b": n15, "c": n17})

    section_divider(
        prs,
        "PART A",
        "Narrative Architecture",
        [
            "From seq 13 — Dynatrace → PagerDuty Four AI Agents",
            "Summary sheet, AS-IS / TO-BE, data flow, agent detail",
            "Enablement, cost, safety, POC afternoon",
        ],
        fill_accent=LIGHT_BLUE,
    )
    append_from(SRC_13, prs, SKIP_13)

    section_divider(
        prs,
        "PART B",
        "Design Diagrams",
        [
            "From seq 15 — Solution Context design pack",
            "AS-IS / TO-BE hub diagrams, Impacted Platforms",
            "Alignment stack, Technical HLD, Agent hub, checklist",
        ],
        fill_accent=GREEN,
    )
    append_from(SRC_15, prs, SKIP_15)

    section_divider(
        prs,
        "PART C",
        "Teams-Only POC Runbook",
        [
            "From seq 17 — detailed Teams-only click path",
            "§0 shared enablement, then SRE → Scribe → Shift → Insights",
            "Honest Slack-first gaps documented for Shift DMs and Insights weekly DMs",
        ],
        fill_accent=ORANGE,
    )
    append_from(SRC_17, prs, SKIP_17)

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        scrub_and_footer(slide, i, total)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote {out}")
    print(f"Total slides: {total} (Part A~{n13}, Part B~{n15}, Part C~{n17} + title/agenda/dividers)")


if __name__ == "__main__":
    base = Path(__file__).resolve().parent
    build(base / "18-Dynatrace-PagerDuty-Four-Agents-Whole-Pack.pptx")
