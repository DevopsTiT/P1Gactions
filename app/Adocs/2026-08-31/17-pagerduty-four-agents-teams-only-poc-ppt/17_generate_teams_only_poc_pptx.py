#!/usr/bin/env python3
"""Detailed PPT from 2026-08-25 seq 11 Teams-only Four Agents POC.
Clean footer — no AXA / ARB template branding."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x2C, 0x5C)
NAVY_MID = RGBColor(0x14, 0x3D, 0x7A)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
GREEN = RGBColor(0x1E, 0x84, 0x4E)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
LIGHT_BLUE = RGBColor(0xD6, 0xE6, 0xF5)
LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
LIGHT_ORANGE = RGBColor(0xFD, 0xE8, 0xD0)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF7)
LIGHT_RED = RGBColor(0xFA, 0xE5, 0xD3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x55, 0x55, 0x55)
YELLOW = RGBColor(0xFF, 0xF2, 0xCC)
PURPLE = RGBColor(0x5B, 0x2C, 0x6F)

FOOTER = "PagerDuty Four Agents — Teams-Only POC"
TOTAL = 20


def set_run(run, text, size=12, bold=False, color=BLACK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_textbox(slide, left, top, width, height, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color)
    return box


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def card(slide, left, top, width, height, title, body, fill=WHITE, title_color=NAVY, line=NAVY_MID, title_size=12, body_size=10):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.06)
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    set_run(r0, title, size=title_size, bold=True, color=title_color)
    if body:
        for i, line in enumerate(body.split("\n")):
            p = tf.add_paragraph()
            r = p.add_run()
            set_run(r, line, size=body_size, color=GRAY)
    return shape


def header(slide, title):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.75), NAVY)
    add_textbox(slide, Inches(0.35), Inches(0.18), Inches(12.5), Inches(0.45), title, size=20, bold=True, color=WHITE)


def key_msg(slide, lines, top=Inches(0.9)):
    h = Inches(0.42 + 0.22 * len(lines))
    add_rect(slide, Inches(0.3), top, Inches(12.7), h, LIGHT_BLUE, NAVY_MID)
    add_textbox(slide, Inches(0.45), top + Inches(0.04), Inches(1.5), Inches(0.22), "Key Message", size=10, bold=True, color=NAVY)
    y = top + Inches(0.26)
    for line in lines:
        add_textbox(slide, Inches(0.45), y, Inches(12.4), Inches(0.2), f"• {line}", size=12, color=BLACK)
        y += Inches(0.22)
    return top + h + Inches(0.08)


def footer(slide, page):
    add_textbox(slide, Inches(0.3), Inches(7.15), Inches(10), Inches(0.25), FOOTER, size=9, color=GRAY)
    add_textbox(slide, Inches(11.2), Inches(7.15), Inches(1.8), Inches(0.25), f"{page} / {TOTAL}", size=10, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_table(slide, left, top, width, rows, col_widths=None, font_size=11):
    nrows = len(rows)
    ncols = len(rows[0])
    row_h = min(0.38, 5.2 / max(nrows, 1))
    table_shape = slide.shapes.add_table(nrows, ncols, left, top, width, Inches(row_h * nrows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            is_header = r_idx == 0
            set_run(run, str(cell_text), size=font_size, bold=is_header, color=WHITE if is_header else BLACK)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = NAVY if is_header else (LIGHT_GRAY if r_idx % 2 else WHITE)
    return table_shape


def bullets(slide, left, top, width, items, size=12, color=BLACK, gap=0.28):
    y = top
    for item in items:
        add_textbox(slide, left, y, width, Inches(gap), f"• {item}", size=size, color=color)
        y += Inches(gap)
    return y


def build(out: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1 Title
    s = blank(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_textbox(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.4), "Detailed runbook deck", size=14, bold=True, color=LIGHT_BLUE)
    add_textbox(s, Inches(0.5), Inches(1.9), Inches(12), Inches(0.7), "PagerDuty Four AI Agents", size=36, bold=True, color=WHITE)
    add_textbox(s, Inches(0.5), Inches(2.7), Inches(12), Inches(0.5), "Teams-Only POC (no Slack)", size=26, bold=True, color=LIGHT_BLUE)
    add_textbox(
        s,
        Inches(0.5),
        Inches(3.4),
        Inches(12),
        Inches(0.4),
        "Source: Daily Files 2026-08-25 / 11-pagerduty-four-agents-teams-only-poc",
        size=13,
        color=WHITE,
    )
    card(s, Inches(0.5), Inches(4.3), Inches(3.8), Inches(1.9), "Full Teams POC", "SRE Agent\nScribe Agent\nInsights conversational Q&A", fill=RGBColor(0x12, 0x3A, 0x6B), title_color=WHITE, line=LIGHT_BLUE)
    card(s, Inches(4.6), Inches(4.3), Inches(3.8), Inches(1.9), "Partial / alternate", "Shift Agent\nPath B: PD web + Calendar\n(Slack DMs not available)", fill=RGBColor(0x12, 0x3A, 0x6B), title_color=WHITE, line=LIGHT_BLUE)
    card(s, Inches(8.7), Inches(4.3), Inches(4.0), Inches(1.9), "Safe demo rule", "Test Service only\nEscalation pages only you\nNever prod on-call", fill=RGBColor(0x12, 0x3A, 0x6B), title_color=WHITE, line=LIGHT_BLUE)

    # 2 Agenda
    s = blank(prs)
    header(s, "Agenda")
    key_msg(s, ["One ordered Teams-only runbook for all four named Advance agents.", "Do shared enablement once, then prove SRE → Scribe → Shift (alt) → Insights."])
    rows = [
        ("#", "Section", "Focus"),
        ("1", "Decision tree + gate", "Which agent; Advance + Admins required"),
        ("2", "§0 Shared enablement", "Teams app, Graph, Advance, linkUser"),
        ("3", "§1 SRE Agent", "Triage / runbook / next steps in Teams"),
        ("4", "§2 Scribe Agent", "Teams meeting transcript + summary"),
        ("5", "§3 Shift Agent", "Honest gap; Path A / Path B"),
        ("6", "§4 Insights Agent", "MTTR Q&A in Teams; weekly DM gap"),
        ("7", "Cheat sheet + schedule", "Cost, surfaces, afternoon plan"),
    ]
    add_table(s, Inches(0.4), Inches(2.0), Inches(12.5), rows, col_widths=[Inches(0.8), Inches(3.5), Inches(8.2)])
    footer(s, 2)

    # 3 Decision tree
    s = blank(prs)
    header(s, "Decision Tree — Which Agent to Prove?")
    y = key_msg(s, ["Pick the agent by the job you need. Shift and Insights weekly DMs have Teams gaps."])
    card(s, Inches(0.35), y, Inches(6.1), Inches(1.15), "Active triage / RCA / runbook?", "→ 1) SRE Agent (Teams Early Access)", fill=LIGHT_BLUE)
    card(s, Inches(6.7), y, Inches(6.2), Inches(1.15), "Bridge notes / PIR draft?", "→ 2) Scribe Agent (Teams meeting + chat)", fill=LIGHT_BLUE)
    y2 = y + Inches(1.3)
    card(s, Inches(0.35), y2, Inches(6.1), Inches(1.55), "On-call OOO / coverage?", "→ 3) Shift Agent\nSlack DMs needed for full path?\n  YES → Path B (web + Calendar)\n  Best-effort → Path A (Advance ask)", fill=LIGHT_ORANGE, title_color=ORANGE, line=ORANGE)
    card(s, Inches(6.7), y2, Inches(6.2), Inches(1.55), "Ops health / MTTR / trends?", "→ 4) Insights Agent (Teams chat)\nWeekly proactive maturity DMs?\n  Slack-only today → document gap", fill=LIGHT_ORANGE, title_color=ORANGE, line=ORANGE)
    card(s, Inches(0.35), y2 + Inches(1.7), Inches(12.55), Inches(1.0), "Always", "TEST service + Level-1 TEST schedule — never prod pages", fill=YELLOW, title_color=ORANGE, line=ORANGE)
    footer(s, 3)

    # 4 Shared gate
    s = blank(prs)
    header(s, "Shared Gate — Do Once Before Any POC")
    key_msg(s, ["Without Advance + both Admins, stop. Then run §0, then POC 1→2→3→4."])
    card(s, Inches(0.4), Inches(2.0), Inches(4.0), Inches(2.2), "Have PagerDuty Advance?", "NO → Sales / trial first\nYES → continue", fill=LIGHT_GREEN, title_color=GREEN, line=GREEN)
    card(s, Inches(4.7), Inches(2.0), Inches(4.0), Inches(2.2), "PD Admin + MS Admin?", "NO → stop; need both\n(Graph + AI Settings)\nYES → §0 Shared enablement", fill=LIGHT_ORANGE, title_color=ORANGE, line=ORANGE)
    card(s, Inches(9.0), Inches(2.0), Inches(3.9), Inches(2.2), "Then run in order", "1 SRE\n2 Scribe\n3 Shift (alt)\n4 Insights", fill=LIGHT_BLUE)
    add_textbox(s, Inches(0.4), Inches(4.5), Inches(12.5), Inches(0.4), "Safe demo rule", size=14, bold=True, color=NAVY)
    bullets(
        s,
        Inches(0.5),
        Inches(5.0),
        Inches(12),
        [
            "Dedicated test Service (example: poc-pd-ai-agents-test)",
            "Escalation on that service must notify only you",
            "Never attach a production primary schedule for Shift POC",
        ],
        size=14,
        gap=0.35,
    )
    footer(s, 4)

    # 5 Takeaway
    s = blank(prs)
    header(s, "Short Takeaway")
    rows = [
        ("Key point", "Detail"),
        ("What this deck is", "Ordered Teams-only runbook for all four named AI agents"),
        ("Builds on", "2026-08-25 seq 9 (agents) + seq 10 (Teams permissions)"),
        ("Full Teams POC", "SRE, Scribe, Insights conversational chat"),
        ("Partial / alternate", "Shift — Slack-first; use Path B or Path A"),
        ("Insights weekly DMs", "Still Slack-first; Teams covers on-demand Q&A only"),
        ("Who spends AI Actions", "SRE and Scribe consume; Shift and Insights usually 0"),
        ("Safe demo rule", "Test Service + escalation that pages only you"),
    ]
    add_table(s, Inches(0.35), Inches(1.1), Inches(12.6), rows, col_widths=[Inches(3.2), Inches(9.4)], font_size=12)
    footer(s, 5)

    # 6 §0 Prerequisites
    s = blank(prs)
    header(s, "§0 Shared Teams Enablement — Prerequisites")
    key_msg(s, ["Do this once. Reuse for all four POCs. Graph scope detail lives in seq 10."])
    rows = [
        ("Need", "What it means", "Why you care"),
        ("PagerDuty Advance", "Paid AI platform + AI Actions budget", "Without it, agent toggles do nothing"),
        ("PD Admin / Owner", "Can open AI Settings and authorize Teams", "Turns agents and chat on"),
        ("Microsoft / Teams Admin", "Consent Graph; allow third-party apps", "Bot cannot read chats or create meetings"),
        ("One standard Teams channel", "Not private or shared", "PD app does not support private/shared"),
        ("Test Service", "e.g. poc-pd-ai-agents-test", "Fake incidents stay off prod rotations"),
        ("Linked users", "Each POC person runs linkUser", "Actions and auto-add need identity map"),
    ]
    add_table(s, Inches(0.3), Inches(1.9), Inches(12.7), rows, col_widths=[Inches(3.0), Inches(5.0), Inches(4.7)], font_size=11)
    footer(s, 6)

    # 7 §0 Steps 1-5
    s = blank(prs)
    header(s, "§0 Enablement Steps (1–5)")
    key_msg(s, ["Admins wire Teams app, Graph consent, and the test Service mapping first."])
    steps = [
        ("1", "Teams Admin", "Allow PagerDuty (or PagerDuty EU) third-party app so users can install it."),
        ("2", "Install", "In Teams: install PagerDuty → Add to a team → pick your one POC team."),
        ("3", "PD Authorize", "Complete the PagerDuty Authorize flow for that team connection."),
        ("4", "Graph consent", "MS Admin accepts Application and/or Delegated scopes (or bot appconnect). Prefer ChatMessage.Read.All (+/or ChatMessage.Read). Prefer User.ReadBasic.All over User.Read.All."),
        ("5", "Map service", "Connect channel ↔ Service poc-pd-ai-agents-test. Escalation must notify only you."),
    ]
    y = Inches(1.85)
    for num, title, body in steps:
        card(s, Inches(0.35), y, Inches(12.6), Inches(0.9), f"{num}. {title}", body, fill=WHITE if int(num) % 2 else LIGHT_GRAY)
        y += Inches(0.95)
    footer(s, 7)

    # 8 §0 Steps 6-10
    s = blank(prs)
    header(s, "§0 Enablement Steps (6–10)")
    key_msg(s, ["Turn Advance chat on, enable agents, link each user, optionally unblock Scribe lobby."])
    steps = [
        ("6", "Advance chat", "PD web → AI → AI Settings → Assistant and AI Agents → Chat Integrations → Teams Connected → toggle Teams On / Enabled."),
        ("7", "Enable agents", "Under AI Agents: SRE, Scribe, Insights = Enabled. Enable Shift only if you want Path A; full Shift DMs need Slack."),
        ("8", "Scope Advance", "Optional but smart: limit Advance access to one PD team that owns the test Service."),
        ("9", "linkUser", "Each POC user: Teams chat with PagerDuty bot → @PagerDuty linkUser. If Delegated-heavy, also run graphAuth."),
        ("10", "Scribe lobby (optional)", "MS Admin may set New-CsApplicationAccessPolicy (US AppId 05ffe668-… or EU 8f79a561-…) and allow Scribe past lobby."),
    ]
    y = Inches(1.85)
    for num, title, body in steps:
        card(s, Inches(0.35), y, Inches(12.6), Inches(0.9), f"{num}. {title}", body, fill=WHITE if int(num) % 2 else LIGHT_GRAY)
        y += Inches(0.95)
    footer(s, 8)

    # 9 §0 Success + Safety
    s = blank(prs)
    header(s, "§0 Success Criteria and Safety")
    add_textbox(s, Inches(0.35), Inches(1.0), Inches(6), Inches(0.3), "Success (pass if)", size=14, bold=True, color=NAVY)
    rows = [
        ("Check", "Pass if"),
        ("App present", "PagerDuty bot responds in the POC team"),
        ("Graph OK", "No UPDATE AVAILABLE blocking Advance features"),
        ("Advance Teams On", "AI Settings shows Teams Connected / Enabled"),
        ("Agents On", "SRE, Scribe, Insights show Enabled"),
        ("Identity", "Your PD user is linked (linkUser done)"),
        ("Safe target", "Test Service maps to POC channel; only you get pages"),
    ]
    add_table(s, Inches(0.35), Inches(1.35), Inches(6.2), rows, col_widths=[Inches(2.0), Inches(4.2)], font_size=10)
    add_textbox(s, Inches(6.9), Inches(1.0), Inches(6), Inches(0.3), "Safety rules", size=14, bold=True, color=NAVY)
    rows2 = [
        ("Rule", "Why"),
        ("One team, one channel, one test Service", "Limits blast radius and Graph surface"),
        ("Standard channel only", "Private/shared channels unsupported"),
        ("No prod escalation", "Prevents waking the real rotation"),
        ("Prefer User.ReadBasic.All", "Least privilege for POC"),
    ]
    add_table(s, Inches(6.9), Inches(1.35), Inches(6.0), rows2, col_widths=[Inches(3.0), Inches(3.0)], font_size=10)
    card(s, Inches(0.35), Inches(5.5), Inches(12.6), Inches(1.2), "AI Actions for §0", "0 for enablement itself. Spend starts when you ask SRE or run Scribe.", fill=LIGHT_GREEN, title_color=GREEN, line=GREEN)
    footer(s, 9)

    # 10 SRE overview
    s = blank(prs)
    header(s, "§1 SRE Agent POC — Goal and Prerequisites")
    key_msg(
        s,
        [
            "Prove SRE summarizes a TEST incident in Teams, uses a small runbook, and suggests next steps.",
            "Human approval required before any remediation. Teams SRE is Early Access.",
        ],
    )
    rows = [
        ("Need", "Detail"),
        ("§0 complete", "Advance Teams Enabled; SRE Enabled; linkUser done"),
        ("Teams Early Access", "Confirm your account has SRE Agent in MS Teams"),
        ("Test Service", "poc-pd-ai-agents-test mapped to the POC channel"),
        ("Runbook file", "One .md or .txt under 100 KB"),
        ("Optional connectors", "Grafana, Datadog, New Relic, CloudWatch, Confluence, GitHub"),
        ("Optional UI", "Incident SRE Agent tab in PD web; Ops Console needs AIOps + Advance"),
    ]
    add_table(s, Inches(0.35), Inches(2.0), Inches(12.6), rows, col_widths=[Inches(3.2), Inches(9.4)], font_size=11)
    footer(s, 10)

    # 11 SRE steps
    s = blank(prs)
    header(s, "§1 SRE Agent — Click / Say Steps")
    steps = [
        "Optional: AI Settings → SRE Agent → add ONE connector you use (Dynatrace creates incidents; usually not the main SRE connector list).",
        "Prepare runbook poc-checkout-latency.md with 5–10 clear steps.",
        "Create a low-risk test incident on poc-pd-ai-agents-test (UI Create incident). No prod routing key.",
        "Open the mapped Teams channel. Confirm the incident card appears.",
        "Type: @pagerduty What are some likely root causes?",
        "Upload / update runbook when offered (or use PD web Incident → SRE Agent tab).",
        "Ask: @pagerduty Analyze past incidents",
        "Ask: @pagerduty What steps should I take first?",
        "If remediation suggested: read it, then approve or decline. Do not auto-run against prod.",
        "Optional: PD web / Ops Console SRE Agent tab.",
        "Resolve the test incident so service memory can save learnings.",
    ]
    y = Inches(0.95)
    for i, t in enumerate(steps, 1):
        add_textbox(s, Inches(0.4), y, Inches(12.5), Inches(0.42), f"{i}. {t}", size=12, color=BLACK)
        y += Inches(0.48)
    footer(s, 11)

    # 12 SRE success / safety / cost
    s = blank(prs)
    header(s, "§1 SRE Agent — Success, Safety, Cost")
    rows = [
        ("Check", "Pass if"),
        ("Triage reply", "Agent posts a grounded summary in Teams (or web tab)"),
        ("Runbook", "A later answer references your uploaded steps"),
        ("Human control", "No remediation ran without an explicit confirm"),
        ("Memory", "After resolve, a later ask on the same service feels more specific"),
    ]
    add_table(s, Inches(0.35), Inches(1.05), Inches(12.6), rows, col_widths=[Inches(2.8), Inches(9.8)], font_size=11)
    rows2 = [
        ("Safety rule", "Why"),
        ("Test Service only", "Avoids waking the real rotation"),
        ("Confirm before remediations", "Agent suggests; humans own blast radius"),
        ("Fact-check log/change claims", "AI can be wrong; verify in monitoring / deploy tools"),
        ("Cap custom details noise", "Only first ~2,000 characters of custom details are analyzed"),
        ("No customer PII in the question", "Teams channel members can see replies"),
    ]
    add_table(s, Inches(0.35), Inches(3.5), Inches(12.6), rows2, col_widths=[Inches(3.8), Inches(8.8)], font_size=11)
    card(s, Inches(0.35), Inches(6.2), Inches(12.6), Inches(0.7), "AI Actions", "4 AI Actions per chat ask or nudge click (also via Incident Workflow / Escalation virtual responder).", fill=YELLOW, title_color=ORANGE, line=ORANGE)
    footer(s, 12)

    # 13 Scribe overview + steps
    s = blank(prs)
    header(s, "§2 Scribe Agent POC — Teams Meeting Path")
    key_msg(
        s,
        [
            "Prove Scribe joins a short Teams meeting on a test incident, streams transcript, posts wrap-up.",
            "Keep POC meeting under ~10 minutes. Human must join within 15 minutes for auto-join.",
        ],
    )
    steps = [
        "Confirm Scribe Agent Enabled (US often needs manual toggle; EU may default on).",
        "Optional: Incident Workflow step Add Scribe Agent on the test Service.",
        "Create test incident on poc-pd-ai-agents-test.",
        "Create/open a short Teams meeting. Copy join URL (include passcode if required).",
        "Paste conference URL onto the incident conference / meeting field in PagerDuty.",
        "Open linked Teams incident channel or chat.",
        "Wait for auto-join OR type: @PagerDuty advance scribe → Add Scribe Agent to meeting → confirm URL.",
        "Join yourself within 15 minutes. Admit Scribe if in lobby.",
        "Speak 2–3 clear sentences: symptom, suspected cause, decision.",
        "End meeting. Confirm transcript activity and post-meeting summary.",
        "Optional: ask Advance for Post-Incident Review draft. Resolve test incident.",
    ]
    y = Inches(1.9)
    for i, t in enumerate(steps, 1):
        add_textbox(s, Inches(0.35), y, Inches(12.6), Inches(0.4), f"{i}. {t}", size=11, color=BLACK)
        y += Inches(0.42)
    footer(s, 13)

    # 14 Scribe success / safety / cost
    s = blank(prs)
    header(s, "§2 Scribe Agent — Success, Safety, Cost")
    rows = [
        ("Check", "Pass if"),
        ("Join", "Scribe appears in the Teams meeting (or joins per product rules)"),
        ("Transcript", "Teams chat or internal capture shows spoken content"),
        ("Summary", "Wrap-up lists decisions, actions, attendees"),
        ("PIR helper", "Later PIR / status draft can reuse bridge context"),
    ]
    add_table(s, Inches(0.35), Inches(1.05), Inches(12.6), rows, col_widths=[Inches(2.5), Inches(10.1)], font_size=11)
    rows2 = [
        ("Safety rule", "Why"),
        ("Test incident + short meeting", "Avoids recording a real customer bridge by accident"),
        ("Warn attendees", "People should know the call is transcribed"),
        ("One Scribe per meeting", "Product limit; do not double-add"),
        ("Cap concurrent meetings", "Up to 10 concurrent Scribe meetings account-wide"),
        ("Keep POC under ~10 minutes", "Limits AI Actions spend"),
    ]
    add_table(s, Inches(0.35), Inches(3.5), Inches(12.6), rows2, col_widths=[Inches(4.0), Inches(8.6)], font_size=11)
    card(s, Inches(0.35), Inches(6.2), Inches(12.6), Inches(0.7), "AI Actions", "~6 per 30 minutes of bridge + ~2 when the final summary posts.", fill=YELLOW, title_color=ORANGE, line=ORANGE)
    footer(s, 14)

    # 15 Shift honest status
    s = blank(prs)
    header(s, "§3 Shift Agent — Honest Teams-Only Status")
    key_msg(s, ["Do not claim full Request coverage → teammate accepts in Teams DM. That path is Slack-first."])
    rows = [
        ("Fact", "What it means"),
        ("Slack-first product", "Docs/pricing describe Shift with Slack; conflict and coverage notifications wire to Slack"),
        ("Teams-only gap", "Full coverage accept automation in Teams DMs is not expected"),
        ("Still give a step path", "Use Path A (best-effort) or Path B (recommended: web + Calendar)"),
    ]
    add_table(s, Inches(0.35), Inches(1.9), Inches(12.6), rows, col_widths=[Inches(3.5), Inches(9.1)], font_size=12)
    card(s, Inches(0.35), Inches(4.3), Inches(6.1), Inches(2.2), "Path A — best-effort", "Ask Advance in Teams about OOO conflict.\nManual override in PD web.\nMark result: partial pass.", fill=LIGHT_ORANGE, title_color=ORANGE, line=ORANGE)
    card(s, Inches(6.7), Inches(4.3), Inches(6.2), Inches(2.2), "Path B — recommended", "Document Slack DM POC deferred.\nGoogle Calendar Extension + OOO block.\nPD web schedule override + history.", fill=LIGHT_GREEN, title_color=GREEN, line=GREEN)
    footer(s, 15)

    # 16 Shift Path B steps
    s = blank(prs)
    header(s, "§3 Shift Path B — Recommended Steps (Teams-Only)")
    key_msg(s, ["Prove OOO vs on-call conflict handling with Calendar + web override. Defer full Slack DMs."])
    steps = [
        "Document: Full Shift Agent Slack DM POC deferred — Teams-only org.",
        "Create test schedule poc-shift-agent with you primary on a known near-term window.",
        "Attach it to Level 1 of a test escalation policy (not prod).",
        "Admin enables Google Calendar Extension; you authorize calendar.",
        "Create an OOO block that overlaps the on-call window.",
        "In PagerDuty web, open the schedule and confirm conflict / overlap is visible.",
        "Create a schedule override for the coverage person for that window.",
        "Confirm the override appears in schedule history.",
        "Optional later: when Slack exists, re-run full Request coverage → accept DM POC.",
    ]
    y = Inches(1.9)
    for i, t in enumerate(steps, 1):
        add_textbox(s, Inches(0.4), y, Inches(12.5), Inches(0.45), f"{i}. {t}", size=13, color=BLACK)
        y += Inches(0.5)
    footer(s, 16)

    # 17 Shift Path A + success
    s = blank(prs)
    header(s, "§3 Shift Path A + Success / Safety / Cost")
    add_textbox(s, Inches(0.35), Inches(0.95), Inches(12), Inches(0.3), "Path A (best-effort) highlights", size=13, bold=True, color=NAVY)
    bullets(
        s,
        Inches(0.45),
        Inches(1.3),
        Inches(12.4),
        [
            "Enable Shift toggle; Level-1 test schedule; Calendar Extension recommended.",
            "Block OOO overlapping your window.",
            "Ask in Teams: I am on vacation <date>. Do I have a conflict?",
            "If Request coverage appears, teammate accept may still need Slack.",
            "Always close with a manual schedule override in PD web.",
        ],
        size=12,
        gap=0.32,
    )
    rows = [
        ("Path", "Pass if"),
        ("Path A", "Conflict language visible in chat OR documented miss; override written manually"),
        ("Path B", "OOO vs on-call handled with Calendar + web override; Slack path explicitly deferred"),
        ("Safety", "No overrides written to a production primary schedule"),
    ]
    add_table(s, Inches(0.35), Inches(3.3), Inches(12.6), rows, col_widths=[Inches(1.8), Inches(10.8)], font_size=11)
    card(s, Inches(0.35), Inches(5.3), Inches(12.6), Inches(1.4), "Safety + cost", "Test schedule / test EP only. Tell the candidate it is a POC. Level-1 only by design. Do not claim Teams DM coverage worked.\nAI Actions = 0 for Shift and for manual web overrides.", fill=LIGHT_GREEN, title_color=GREEN, line=GREEN)
    footer(s, 17)

    # 18 Insights
    s = blank(prs)
    header(s, "§4 Insights Agent — Conversational in Teams")
    key_msg(
        s,
        [
            "Prove Insights answers MTTR / MTTA / volume questions inside Microsoft Teams.",
            "Weekly proactive maturity DMs remain Slack-first — document that gap.",
        ],
    )
    steps = [
        "Confirm Insights Agent Enabled.",
        "Pick one Team/Service you own with real history (sanity-checkable).",
        "Open POC Teams channel (or Advance chat).",
        "Ask: @pagerduty How many high urgency incidents were there last week on <ServiceOrTeam>?",
        "Ask: @pagerduty How has the average time to resolve changed over the past 6 complete months for <Team>?",
        "Ask: @pagerduty Was MTTA faster this month than last month for <Service>?",
        "Compare each answer to Analytics in the PagerDuty web UI.",
        "Optional: Rate AI Response. Document weekly DM Slack gap. Apply any config tips on TEST first.",
    ]
    y = Inches(1.9)
    for i, t in enumerate(steps, 1):
        add_textbox(s, Inches(0.35), y, Inches(12.6), Inches(0.45), f"{i}. {t}", size=12, color=BLACK)
        y += Inches(0.5)
    footer(s, 18)

    # 19 Insights success + cheat sheet
    s = blank(prs)
    header(s, "§4 Insights Success + Cost / Surface Cheat Sheet")
    rows = [
        ("Check", "Pass if"),
        ("Chat answer", "Number or trend tied to your Team/Service in Teams"),
        ("Sanity", "Roughly matches Analytics UI for the same filter"),
        ("Visibility", "Channel members can see the @pagerduty reply"),
        ("DM gap documented", "Notes state weekly proactive DMs are Slack-first"),
    ]
    add_table(s, Inches(0.3), Inches(1.0), Inches(12.7), rows, col_widths=[Inches(3.0), Inches(9.7)], font_size=11)
    rows2 = [
        ("Agent", "Teams-only demo surface", "AI Actions", "Honest note"),
        ("SRE", "Teams @pagerduty ask; optional PD web tab", "4 per ask", "Early Access in Teams"),
        ("Scribe", "Teams meeting URL + advance scribe", "~6/30m + ~2 summary", "Strongest Teams-only fit"),
        ("Shift", "Path A ask; Path B web + Calendar", "0", "Full coverage DMs need Slack"),
        ("Insights", "Teams @pagerduty analytics Q&A", "0", "Weekly maturity DMs Slack-first"),
    ]
    add_table(s, Inches(0.3), Inches(3.7), Inches(12.7), rows2, col_widths=[Inches(1.6), Inches(5.0), Inches(2.4), Inches(3.7)], font_size=11)
    footer(s, 19)

    # 20 Run order + data flow + close
    s = blank(prs)
    header(s, "Afternoon Plan + Data Flow + Commands")
    rows = [
        ("Order", "Block", "Time box"),
        ("0", "Shared Teams enablement", "30–60 min (Admins)"),
        ("1", "SRE Agent POC", "15–20 min"),
        ("2", "Scribe Agent POC", "15–20 min (meeting ≤10 min)"),
        ("3", "Shift Path B (or Path A)", "10–15 min"),
        ("4", "Insights conversational", "10–15 min"),
    ]
    add_table(s, Inches(0.3), Inches(0.95), Inches(7.0), rows, col_widths=[Inches(1.0), Inches(3.5), Inches(2.5)], font_size=11)
    card(
        s,
        Inches(7.5),
        Inches(0.95),
        Inches(5.4),
        Inches(3.0),
        "Data flow",
        "§0 Graph + Authorize + linkUser\n  → 1 SRE triage in Teams\n  → 2 Scribe on Teams meeting\n  → 3 Shift Path B (web/Calendar)\n  → 4 Insights Q&A in Teams\nWeekly DMs = Slack gap",
        fill=LIGHT_BLUE,
        body_size=11,
    )
    add_textbox(s, Inches(0.3), Inches(4.2), Inches(12), Inches(0.3), "Type in Teams (not shell)", size=13, bold=True, color=NAVY)
    bullets(
        s,
        Inches(0.4),
        Inches(4.55),
        Inches(12.4),
        [
            "@PagerDuty linkUser",
            "graphAuth   (if Delegated-heavy)",
            "@pagerduty What are some likely root causes?",
            "@PagerDuty advance scribe",
            "@pagerduty How has the average time to resolve changed over the past 6 complete months for <Team>?",
        ],
        size=12,
        gap=0.35,
    )
    footer(s, 20)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote {out} ({TOTAL} slides)")


if __name__ == "__main__":
    base = Path(__file__).resolve().parent
    build(base / "17-PagerDuty-Four-Agents-Teams-Only-POC.pptx")
