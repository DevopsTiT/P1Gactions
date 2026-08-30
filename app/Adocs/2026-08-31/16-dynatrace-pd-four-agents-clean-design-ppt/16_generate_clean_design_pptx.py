#!/usr/bin/env python3
"""Clean design PPT: Dynatrace → PagerDuty Four AI Agents.
No AXA / ARB template branding or footer text."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x2C, 0x5C)
NAVY_MID = RGBColor(0x14, 0x3D, 0x7A)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
GREEN = RGBColor(0x1E, 0x84, 0x4E)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
LIGHT_BLUE = RGBColor(0xD6, 0xE6, 0xF5)
LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
LIGHT_ORANGE = RGBColor(0xFD, 0xE8, 0xD0)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF7)
LIGHT_RED = RGBColor(0xFA, 0xE5, 0xD3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x55, 0x55, 0x55)
YELLOW = RGBColor(0xFF, 0xF2, 0xCC)


def set_run(run, text, size=12, bold=False, color=BLACK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_textbox(slide, left, top, width, height, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color)
    return box


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def card(slide, left, top, width, height, title, body, fill=WHITE, title_color=NAVY, line=NAVY_MID):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.08)
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    set_run(r0, title, size=12, bold=True, color=title_color)
    if body:
        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        set_run(r1, body, size=10, color=GRAY)
    return shape


def arrow(slide, left, top, width=Inches(0.45)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.22))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY_MID
    shape.line.fill.background()
    return shape


def header(slide, title):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.8), NAVY)
    add_textbox(slide, Inches(0.35), Inches(0.18), Inches(12.5), Inches(0.45), title, size=22, bold=True, color=WHITE)


def key_msg(slide, lines):
    add_rect(slide, Inches(0.3), Inches(0.95), Inches(12.7), Inches(0.55 + 0.22 * (len(lines) - 1)), LIGHT_BLUE, NAVY_MID)
    add_textbox(slide, Inches(0.45), Inches(1.0), Inches(1.5), Inches(0.25), "Key Message", size=11, bold=True, color=NAVY)
    y = Inches(1.22)
    for line in lines:
        add_textbox(slide, Inches(0.45), y, Inches(12.4), Inches(0.22), f"• {line}", size=12, color=BLACK)
        y += Inches(0.22)


def footer(slide, page, total=5):
    """Minimal footer only — deck name + page. No company or board-template labels."""
    add_textbox(
        slide,
        Inches(0.3),
        Inches(7.15),
        Inches(10),
        Inches(0.25),
        "Dynatrace to PagerDuty Four AI Agents — Design",
        size=9,
        color=GRAY,
    )
    add_textbox(
        slide,
        Inches(11.2),
        Inches(7.15),
        Inches(1.8),
        Inches(0.25),
        f"{page} / {total}",
        size=10,
        bold=True,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
    )


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def badge(slide, left, top, text, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.5), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text, size=14, bold=True, color=WHITE)
    return shape


def build(out: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = 5

    # ----- 1 Title -----
    s = blank(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_textbox(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.4), "Architecture Design Pack", size=16, bold=True, color=LIGHT_BLUE)
    add_textbox(s, Inches(0.5), Inches(1.9), Inches(12), Inches(0.7), "Dynatrace to PagerDuty", size=36, bold=True, color=WHITE)
    add_textbox(s, Inches(0.5), Inches(2.7), Inches(12), Inches(0.5), "Four AI Agents — Design Diagrams", size=24, bold=True, color=LIGHT_BLUE)
    add_textbox(
        s,
        Inches(0.5),
        Inches(3.4),
        Inches(12),
        Inches(0.4),
        "SRE  |  Scribe  |  Shift  |  Insights",
        size=16,
        color=WHITE,
    )
    card(
        s,
        Inches(0.5),
        Inches(4.3),
        Inches(5.8),
        Inches(1.8),
        "Purpose",
        "Design-review diagrams for the assist layer\non Dynatrace-fed PagerDuty incidents.\nClean deck — no board-template branding.",
        fill=RGBColor(0x12, 0x3A, 0x6B),
        title_color=WHITE,
        line=LIGHT_BLUE,
    )
    card(
        s,
        Inches(6.6),
        Inches(4.3),
        Inches(5.8),
        Inches(1.8),
        "Contents",
        "1 Title\n2 AS-IS Solution Context\n3 TO-BE Solution Context\n4 Technical HLD\n5 Agent Hub Detail",
        fill=RGBColor(0x12, 0x3A, 0x6B),
        title_color=WHITE,
        line=LIGHT_BLUE,
    )

    # ----- 2 AS-IS -----
    s = blank(prs)
    header(s, "Solution Context Diagram — AS-IS")
    key_msg(
        s,
        [
            "Dynatrace detects problems and notifies PagerDuty; escalation pages on-call.",
            "Triage, bridge notes, coverage, and analytics are still mostly manual human work.",
        ],
    )
    badge(s, Inches(0.35), Inches(1.85), "AS-IS", ORANGE)

    card(s, Inches(0.4), Inches(2.5), Inches(2.6), Inches(1.25), "Dynatrace", "Davis Problem\nmetrics / logs / traces", fill=LIGHT_ORANGE, title_color=ORANGE)
    arrow(s, Inches(3.1), Inches(3.0))
    card(s, Inches(3.65), Inches(2.5), Inches(2.7), Inches(1.25), "PagerDuty", "Incident created\nEscalation Policy", fill=LIGHT_BLUE)
    arrow(s, Inches(6.45), Inches(3.0))
    card(
        s,
        Inches(7.0),
        Inches(2.4),
        Inches(2.9),
        Inches(1.45),
        "Human on-call",
        "Ack · dig docs · type notes\nAsk coverage · export CSV",
        fill=YELLOW,
        title_color=ORANGE,
        line=ORANGE,
    )
    card(
        s,
        Inches(10.2),
        Inches(2.4),
        Inches(2.8),
        Inches(1.45),
        "Pain",
        "Slow first 10 min\nLost bridge decisions\nCoverage gaps\nManual trends",
        fill=LIGHT_RED,
        title_color=ACCENT_RED,
        line=ACCENT_RED,
    )

    card(s, Inches(0.4), Inches(4.2), Inches(2.6), Inches(1.2), "Apps / Infra", "Checkout · hosts · DB\n(monitored by Dynatrace)", fill=WHITE)
    card(s, Inches(3.3), Inches(4.2), Inches(2.8), Inches(1.2), "Chat (optional)", "Teams / Slack cards only\nNo AI agents yet", fill=WHITE)
    card(s, Inches(6.4), Inches(4.2), Inches(3.0), Inches(1.2), "Schedule", "Level-1 on-call\nManual OOO overrides", fill=WHITE)
    card(s, Inches(9.7), Inches(4.2), Inches(3.2), Inches(1.2), "Analytics", "PD Analytics UI\nManual pull for leadership", fill=WHITE)

    add_textbox(
        s,
        Inches(0.4),
        Inches(5.7),
        Inches(12.5),
        Inches(0.9),
        "AS-IS gap: Detect (Dynatrace) and Page (PagerDuty) exist. The assist layer between page and resolve is human-only.",
        size=13,
        bold=True,
        color=NAVY,
    )
    footer(s, 2, total)

    # ----- 3 TO-BE -----
    s = blank(prs)
    header(s, "Solution Context Diagram — TO-BE")
    key_msg(
        s,
        [
            "PagerDuty Advance Four AI Agents assist beside the human after Dynatrace creates the incident.",
            "Optional Event Orchestration cuts noise before agents run. Humans still confirm remediations.",
        ],
    )
    badge(s, Inches(0.35), Inches(1.85), "TO-BE", GREEN)

    card(s, Inches(0.3), Inches(2.45), Inches(2.4), Inches(1.1), "Dynatrace", "Problem + problem URL\nin PD payload", fill=LIGHT_GREEN, title_color=GREEN)
    card(s, Inches(0.3), Inches(3.7), Inches(2.4), Inches(1.0), "Optional EO", "Route / pause / group\n(noise reduction)", fill=LIGHT_GRAY)
    card(s, Inches(0.3), Inches(4.9), Inches(2.4), Inches(1.0), "Ops", "Teams / PD web\nSSO access as configured", fill=WHITE)

    arrow(s, Inches(2.8), Inches(2.9))
    arrow(s, Inches(2.8), Inches(4.05))
    arrow(s, Inches(2.8), Inches(5.25))

    hub = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.4), Inches(2.7), Inches(3.3), Inches(2.6))
    hub.fill.solid()
    hub.fill.fore_color.rgb = WHITE
    hub.line.color.rgb = ACCENT_RED
    hub.line.width = Pt(2.5)
    tf = hub.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, "PagerDuty", size=14, bold=True, color=ACCENT_RED)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    set_run(r2, "Incident hub", size=11, color=GRAY)
    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    set_run(r3, "\nAdvance AI Agents", size=13, bold=True, color=NAVY)
    p4 = tf.add_paragraph()
    r4 = p4.add_run()
    set_run(r4, "SRE · Scribe\nShift · Insights", size=12, bold=True, color=TEAL)

    arrow(s, Inches(6.85), Inches(3.8))

    card(s, Inches(7.4), Inches(2.4), Inches(2.7), Inches(1.05), "SRE Agent", "Triage + runbook\n4 Actions / ask", fill=LIGHT_BLUE)
    card(s, Inches(10.3), Inches(2.4), Inches(2.7), Inches(1.05), "Scribe Agent", "Meeting transcript\n~6/30m + 2", fill=LIGHT_BLUE)
    card(s, Inches(7.4), Inches(3.6), Inches(2.7), Inches(1.05), "Shift Agent", "Coverage / override\n0 Actions", fill=LIGHT_BLUE)
    card(s, Inches(10.3), Inches(3.6), Inches(2.7), Inches(1.05), "Insights Agent", "MTTR Q&A\n0 Actions", fill=LIGHT_BLUE)
    card(
        s,
        Inches(7.4),
        Inches(4.9),
        Inches(5.6),
        Inches(1.15),
        "Human decides",
        "Confirm remediations · Own blast radius · Resolve → SRE memory",
        fill=YELLOW,
        title_color=ORANGE,
        line=ORANGE,
    )
    footer(s, 3, total)

    # ----- 4 HLD -----
    s = blank(prs)
    header(s, "Technical Infrastructure — HLD Overview")
    key_msg(
        s,
        [
            "Dynatrace sends problems to PagerDuty; Advance agents attach to the incident and chat/meeting surfaces.",
            "Teams path needs Graph message-read for Advance; Shift full DMs prefer Slack.",
        ],
    )

    card(s, Inches(0.3), Inches(2.2), Inches(2.5), Inches(1.3), "Responder", "Teams / Slack\nlinkUser · graphAuth*", fill=WHITE)
    card(s, Inches(0.3), Inches(3.7), Inches(2.5), Inches(1.3), "Admin", "PD AI Settings\nMS Graph consent", fill=WHITE)
    card(s, Inches(0.3), Inches(5.2), Inches(2.5), Inches(1.1), "Calendar", "Google Calendar Ext.\n(Shift)", fill=LIGHT_GRAY)

    arrow(s, Inches(2.9), Inches(2.75))
    arrow(s, Inches(2.9), Inches(4.25))

    card(
        s,
        Inches(3.5),
        Inches(2.2),
        Inches(3.2),
        Inches(2.0),
        "PagerDuty",
        "Service · Incident\nAdvance Agents\nAnalytics",
        fill=LIGHT_BLUE,
        title_color=NAVY,
        line=ACCENT_RED,
    )
    card(s, Inches(3.5), Inches(4.4), Inches(3.2), Inches(1.5), "Optional EO", "Event Orchestration\nAlert Grouping", fill=LIGHT_GRAY)

    arrow(s, Inches(6.85), Inches(3.0))

    card(
        s,
        Inches(7.4),
        Inches(2.2),
        Inches(2.7),
        Inches(1.5),
        "Dynatrace",
        "Problem / Davis\nNotify → PD key\nproblem URL",
        fill=LIGHT_GREEN,
        title_color=GREEN,
        line=GREEN,
    )
    card(s, Inches(7.4), Inches(3.9), Inches(2.7), Inches(1.3), "Apps / Infra", "Checkout · hosts\nDB · JVM (monitored)", fill=WHITE)

    card(s, Inches(10.3), Inches(2.2), Inches(2.7), Inches(1.2), "Teams meeting", "Scribe join\ntranscript", fill=LIGHT_BLUE)
    card(s, Inches(10.3), Inches(3.55), Inches(2.7), Inches(1.2), "Teams / Slack chat", "@pagerduty asks\nSRE · Insights", fill=LIGHT_BLUE)
    card(
        s,
        Inches(10.3),
        Inches(4.9),
        Inches(2.7),
        Inches(1.2),
        "PD web",
        "SRE Agent tab\nSchedules override",
        fill=LIGHT_ORANGE,
        title_color=ORANGE,
        line=ORANGE,
    )

    add_textbox(
        s,
        Inches(0.3),
        Inches(6.5),
        Inches(12.7),
        Inches(0.4),
        "* graphAuth only if tenant is Delegated-heavy. Prefer Application ChatMessage.Read.All for Advance.",
        size=11,
        color=GRAY,
    )
    footer(s, 4, total)

    # ----- 5 Agent hub -----
    s = blank(prs)
    header(s, "Agent Hub Detail — Four Agents on One Incident")
    key_msg(s, ["One Dynatrace-fed incident can use multiple agents; pick by job."])

    card(
        s,
        Inches(4.9),
        Inches(2.5),
        Inches(3.5),
        Inches(1.6),
        "PD Incident",
        "from Dynatrace Problem\n+ custom_details / problem URL",
        fill=YELLOW,
        title_color=ORANGE,
        line=ORANGE,
    )

    card(
        s,
        Inches(0.4),
        Inches(2.3),
        Inches(3.8),
        Inches(1.9),
        "SRE Agent",
        "Summarize · runbook · past incidents\nNext steps · human confirm\nCost: 4 Actions / ask\nTeams EA + PD web tab",
        fill=LIGHT_BLUE,
    )
    card(
        s,
        Inches(8.9),
        Inches(2.3),
        Inches(4.0),
        Inches(1.9),
        "Scribe Agent",
        "Join Teams/Zoom/Meet\nTranscript + wrap-up → PIR\nCost: ~6/30m + 2\nHuman join ≤15 min",
        fill=LIGHT_BLUE,
    )
    card(
        s,
        Inches(0.4),
        Inches(4.5),
        Inches(3.8),
        Inches(1.9),
        "Shift Agent",
        "Level-1 OOO conflict\nSlack DM coverage OR\nPath B: Calendar + web override\nCost: 0 Actions",
        fill=LIGHT_GREEN,
        title_color=GREEN,
        line=GREEN,
    )
    card(
        s,
        Inches(8.9),
        Inches(4.5),
        Inches(4.0),
        Inches(1.9),
        "Insights Agent",
        "MTTR / MTTA / volume Q&A\nSanity-check Analytics UI\nWeekly tips: Slack-first\nCost: 0 Actions",
        fill=LIGHT_GREEN,
        title_color=GREEN,
        line=GREEN,
    )
    card(
        s,
        Inches(4.9),
        Inches(4.5),
        Inches(3.5),
        Inches(1.9),
        "Human",
        "Owns blast radius\nApprove remediations\nResolve → SRE memory",
        fill=WHITE,
        line=NAVY,
    )
    footer(s, 5, total)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote {out}")


if __name__ == "__main__":
    base = Path(__file__).resolve().parent
    build(base / "16-Dynatrace-PagerDuty-Four-Agents-Clean-Design.pptx")
