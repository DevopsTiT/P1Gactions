#!/usr/bin/env python3
"""Generate AXA-Japan-ARB-style PPT for P260120F Splunk to Dynatrace migration."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# Widescreen 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x2C, 0x5C)
NAVY_MID = RGBColor(0x14, 0x3D, 0x7A)
BLUE_HDR = RGBColor(0x1F, 0x4E, 0x8C)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
GREEN = RGBColor(0x1E, 0x84, 0x4E)
LIGHT_BLUE = RGBColor(0xD6, 0xE6, 0xF5)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x55, 0x55, 0x55)
YELLOW = RGBColor(0xFF, 0xF2, 0xCC)


def set_run(run, text, size=12, bold=False, color=BLACK, font_name="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def add_textbox(slide, left, top, width, height, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color)
    return box


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def add_round_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def footer(slide, page, total=12):
    add_textbox(
        slide,
        Inches(0.3),
        Inches(7.15),
        Inches(8),
        Inches(0.3),
        "AXA Japan ARB Template style · P260120F Splunk to Dynatrace",
        size=9,
        color=GRAY,
    )
    add_textbox(
        slide,
        Inches(11.5),
        Inches(7.15),
        Inches(1.5),
        Inches(0.3),
        f"{page}",
        size=10,
        bold=True,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
    )


def section_header_bar(slide, title):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), NAVY)
    add_textbox(slide, Inches(0.4), Inches(0.2), Inches(12), Inches(0.5), title, size=24, bold=True, color=WHITE)


def key_message_box(slide, lines, top=Inches(1.0)):
    add_rect(slide, Inches(0.35), top, Inches(12.6), Inches(0.7 + 0.22 * max(0, len(lines) - 1)), LIGHT_BLUE, NAVY_MID)
    add_textbox(slide, Inches(0.5), top + Inches(0.08), Inches(1.4), Inches(0.3), "Key Message", size=11, bold=True, color=NAVY)
    y = top + Inches(0.32)
    for line in lines:
        add_textbox(slide, Inches(0.5), y, Inches(12.2), Inches(0.28), f"• {line}", size=12, color=BLACK)
        y += Inches(0.25)
    return y + Inches(0.15)


def add_table(slide, left, top, width, rows, col_widths=None, header=True):
    cols = len(rows[0])
    table_shape = slide.shapes.add_table(len(rows), cols, left, top, width, Inches(0.32 * len(rows)))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            is_hdr = header and r_idx == 0
            set_run(run, str(cell_text), size=10 if not is_hdr else 11, bold=is_hdr, color=WHITE if is_hdr else BLACK)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE_HDR if is_hdr else (LIGHT_GRAY if r_idx % 2 else WHITE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape


def flow_box(slide, left, top, width, height, title, body, fill=WHITE, title_color=NAVY):
    shape = add_round_rect(slide, left, top, width, height, fill, NAVY_MID)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    set_run(r0, title, size=11, bold=True, color=title_color)
    p1 = tf.add_paragraph()
    r1 = p1.add_run()
    set_run(r1, body, size=9, color=GRAY)
    return shape


def arrow_right(slide, left, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, Inches(0.45), Inches(0.25))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY_MID
    shape.line.fill.background()
    return shape


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def build(out_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ----- 1 Title -----
    s = blank_slide(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_textbox(s, Inches(0.6), Inches(0.5), Inches(6), Inches(0.4), "Architecture Review Board", size=16, bold=True, color=LIGHT_BLUE)
    add_textbox(s, Inches(0.6), Inches(1.4), Inches(12), Inches(1.0), "Architecture Design Proposal", size=36, bold=True, color=WHITE)
    add_textbox(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5), "for (P260120F - Splunk to Dynatrace migration)", size=20, color=LIGHT_BLUE)

    meta = [
        ("Objective", "Build Permit + Cloud Permit"),
        ("Gate", "Gate2"),
        ("Architecture Governance", "Standard"),
        ("Presenter (Team Name)", "ROA Rosel (Ops Middleware and Monitoring)"),
        ("Line of Business", "OneAXA"),
        ("ARB Date", "2026/06/26"),
        ("Engineering Point of Contact", "YASUDA Masayuki, HUNG Hsinyi, BARBASTE Nicolas"),
        ("Ops Point of Contact", "YASUDA Masayuki"),
        ("Platform Architecture Point of Contact", "ROA Rosel"),
    ]
    add_table(
        s,
        Inches(0.6),
        Inches(3.2),
        Inches(8.5),
        [("Field", "Value")] + meta,
        col_widths=[Inches(3.2), Inches(5.3)],
    )
    add_textbox(s, Inches(0.6), Inches(6.9), Inches(10), Inches(0.3), "Architecture / AXA Japan (template style) · LAST UPDATED 2026/08/31", size=10, color=LIGHT_BLUE)

    # ----- 2 Scope / Approval matrix (simplified) -----
    s = blank_slide(prs)
    section_header_bar(s, "This Approval is requested for")
    add_textbox(s, Inches(8.5), Inches(0.25), Inches(4.5), Inches(0.4), "Here's the Proposal Scope", size=12, bold=True, color=ACCENT_RED)
    add_textbox(
        s,
        Inches(0.4),
        Inches(1.05),
        Inches(12.5),
        Inches(0.4),
        "Red box = this proposal scope (Gate2 Build Permit + Cloud Permit). Matrix simplified from ARB template.",
        size=12,
        color=GRAY,
    )
    rows = [
        ("#", "Slide Title / Required for", "PA", "BP", "CP", "IR"),
        ("1", "Architecture Design <Summary Sheet>", "M", "M", "M", "M"),
        ("2", "Updates", "M", "M", "", "M"),
        ("3", "Impacted Applications", "M", "M", "", ""),
        ("4", "Solution Context Diagram", "M", "M", "", ""),
        ("5", "Impacted Platforms", "M", "M", "M", ""),
        ("6", "Alignment with DPRA", "M", "M", "", ""),
        ("8", "Data Architecture", "", "M", "M", ""),
        ("10", "Technical/Infrastructure Architecture -> HLD", "", "M*2", "M", ""),
        ("12", "Security Architecture", "", "M*2", "M*2", "M"),
        ("15", "Cost Estimate", "", "M", "", ""),
        ("16", "Decommission Plan", "", "M", "M", ""),
        ("18", "Cloud Permit - Cloud Binding Checklist", "n/a", "n/a", "M", "n/a"),
    ]
    add_table(s, Inches(0.35), Inches(1.5), Inches(12.6), rows, col_widths=[Inches(0.5), Inches(5.5), Inches(1.3), Inches(1.5), Inches(1.5), Inches(1.3)])
    add_textbox(s, Inches(0.4), Inches(6.5), Inches(12), Inches(0.4), "M = Mandatory · M*2 = Mandatory when template requires · AI Governance N/A (not an AI solution)", size=11, color=GRAY)
    footer(s, 2)

    # ----- 3 Summary Sheet -----
    s = blank_slide(prs)
    section_header_bar(s, "Architecture Design <Summary Sheet>")
    summary_rows = [
        ("Item", "Answer"),
        ("Objective <Gate>", "Solution Approval <Gate2>"),
        ("Archi Gov.", "Standard"),
        (
            "Background WHY",
            "Decommission Splunk; migrate to AXA-GO Dynatrace. Cut license/maintenance cost; reduce MPI footprint; gain observability; enterprise alignment.",
        ),
        (
            "Architecture Design Overview",
            "Provision AXA-GO Dynatrace SaaS dedicated tenant. Instrument apps to Dynatrace logging. Migrate Splunk monitoring/alerting. Add Ex-ADL. Decommission Splunk MPI.",
        ),
        ("Alignment", "Aligned with D-PRA / Architecture Roadmap"),
        ("API", "N/A"),
        ("Cloud", "Dynatrace"),
        ("Project Cost", "25 (M Yen / Gate 2)"),
        ("Production Release Timing", "Q3-2026"),
        ("Security Review Status", "TBA – before Friday (update before ARB)"),
        (
            "Architect View (in charge)",
            "Support migration to reduce MPI servers and centralize logs to Dynatrace SaaS aligned with AXA-GO roadmap.",
        ),
    ]
    add_table(s, Inches(0.35), Inches(1.1), Inches(12.6), summary_rows, col_widths=[Inches(3.0), Inches(9.6)])
    footer(s, 3)

    # ----- 4 Update / Approval History -----
    s = blank_slide(prs)
    section_header_bar(s, "Update and Approval History")
    key_message_box(s, ["For build/cloud permit approval"])
    add_textbox(s, Inches(0.4), Inches(2.0), Inches(6), Inches(0.3), "Technical updates", size=14, bold=True, color=NAVY)
    add_table(
        s,
        Inches(0.4),
        Inches(2.4),
        Inches(12.5),
        [
            ("#", "Item", "Update detail"),
            ("1", "Technical and Interface", "Added detailed HLD for the migration"),
            ("2", "", ""),
            ("3", "", ""),
        ],
        col_widths=[Inches(0.6), Inches(3.5), Inches(8.4)],
    )
    add_textbox(s, Inches(0.4), Inches(4.2), Inches(6), Inches(0.3), "ARB Approval History", size=14, bold=True, color=NAVY)
    add_table(
        s,
        Inches(0.4),
        Inches(4.6),
        Inches(12.5),
        [
            ("#", "ARB Approval History (Date)", "ARB Page Link"),
            ("1", "PA (2026/05/15)", "OneAXA ARB Splunk to Dynatrace – Confluence by AXA GO"),
            ("2", "BP (2026/06/26)", "OneAXA ARB P260120F – THIS ARB"),
            ("3", "IR (YYYY/MM/DD)", ""),
            ("4", "CP (YYYY/MM/DD)", ""),
        ],
        col_widths=[Inches(0.6), Inches(3.5), Inches(8.4)],
    )
    footer(s, 4)

    # ----- 5 AS-IS -----
    s = blank_slide(prs)
    section_header_bar(s, "Solution Context Diagram — AS-IS")
    key_message_box(
        s,
        [
            "MPI platform hosts Splunk for centralized application logging (ADJ + ALJ).",
            "OpenPaaS and Lambda stream logs to AWS S3; Splunk ingests for analysis.",
        ],
    )
    add_round_rect(s, Inches(0.4), Inches(2.3), Inches(1.6), Inches(0.55), ORANGE)
    add_textbox(s, Inches(0.55), Inches(2.4), Inches(1.3), Inches(0.4), "AS-IS", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    flow_box(s, Inches(0.5), Inches(3.2), Inches(2.6), Inches(1.3), "AXA-GO OpenPaaS", "Multi Log Forwarder\n→ S3 bucket")
    arrow_right(s, Inches(3.2), Inches(3.7))
    flow_box(s, Inches(3.75), Inches(3.2), Inches(2.6), Inches(1.3), "ALJ/ADJ AWS S3", "Logs from Lambda\nand OpenPaaS")
    arrow_right(s, Inches(6.45), Inches(3.7))
    hub = flow_box(s, Inches(7.0), Inches(3.0), Inches(2.8), Inches(1.7), "ALJ MPI — Splunk", "Central ingest &\nsearch (hub)", fill=YELLOW)
    flow_box(s, Inches(10.2), Inches(2.5), Inches(2.7), Inches(1.1), "ADJ MPI Apps", "Applications → Splunk")
    flow_box(s, Inches(10.2), Inches(3.8), Inches(2.7), Inches(1.1), "ALJ POD / Local HUB", "Applications → Splunk")
    flow_box(s, Inches(0.5), Inches(5.0), Inches(2.6), Inches(1.2), "AXA HQ Ops", "Shirokane operations\nusers → Splunk")
    arrow_right(s, Inches(3.2), Inches(5.5))
    add_textbox(s, Inches(3.75), Inches(5.4), Inches(5), Inches(0.4), "→ Splunk hub (center of AS-IS)", size=12, bold=True, color=NAVY)
    footer(s, 5)

    # ----- 6 TO-BE -----
    s = blank_slide(prs)
    section_header_bar(s, "Solution Context Diagram — TO-BE")
    key_message_box(
        s,
        [
            "Dynatrace (SaaS) from AXA-GO will replace Splunk.",
            "Ex-ADL servers are added in scope for Dynatrace application logs.",
        ],
    )
    add_round_rect(s, Inches(0.4), Inches(2.3), Inches(1.6), Inches(0.55), GREEN)
    add_textbox(s, Inches(0.55), Inches(2.4), Inches(1.3), Inches(0.4), "TO-BE", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    flow_box(s, Inches(0.4), Inches(3.1), Inches(2.5), Inches(1.15), "OpenPaaS", "Multi Log Forwarder\n→ Dynatrace direct")
    flow_box(s, Inches(0.4), Inches(4.4), Inches(2.5), Inches(1.15), "AWS Firehose", "Kinesis Data Firehose\n→ Dynatrace")
    flow_box(s, Inches(0.4), Inches(5.7), Inches(2.5), Inches(0.95), "AXA HQ Ops", "Console via proxy\n+ OneAccount")

    hub = add_round_rect(s, Inches(5.2), Inches(3.6), Inches(3.2), Inches(1.8), WHITE, ACCENT_RED)
    tf = hub.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, "Dynatrace SaaS", size=16, bold=True, color=ACCENT_RED)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    set_run(r2, "AXA-GO / OneAXA hub\nIngest · Store · Query", size=11, color=GRAY)

    flow_box(s, Inches(9.5), Inches(2.9), Inches(3.3), Inches(1.0), "ADJ MPI Apps", "Instrumented → Dynatrace")
    flow_box(s, Inches(9.5), Inches(4.1), Inches(3.3), Inches(1.0), "ALJ POD Apps", "→ Dynatrace")
    flow_box(s, Inches(9.5), Inches(5.3), Inches(3.3), Inches(1.1), "Ex-ADL / Local HUB", "NEW in scope → Dynatrace", fill=LIGHT_BLUE)

    add_textbox(s, Inches(3.2), Inches(3.9), Inches(1.8), Inches(0.3), "────────►", size=14, bold=True, color=NAVY)
    add_textbox(s, Inches(8.5), Inches(4.2), Inches(1.0), Inches(0.3), "◄────", size=14, bold=True, color=NAVY)
    footer(s, 6)

    # ----- 7 Impacted Platforms -----
    s = blank_slide(prs)
    section_header_bar(s, "Impacted Platforms")
    key_message_box(
        s,
        ["Dynatrace will consolidate logging from system to application and speed up troubleshooting during incidents."],
    )
    add_textbox(s, Inches(0.4), Inches(2.1), Inches(12), Inches(0.3), "Highlighted impact (from ARB map)", size=14, bold=True, color=NAVY)
    add_table(
        s,
        Inches(0.4),
        Inches(2.5),
        Inches(12.5),
        [
            ("Area", "Impact", "Note"),
            ("IT / IT Platform (IT4IT)", "Primary", "Shared logging/monitoring platform change"),
            ("Cloud (Infrastructure & Cloud)", "Primary", "SaaS Dynatrace + AWS Firehose/S3 paths"),
            ("Business domains (green on map)", "Reuse AXA-GO", "Consumers of shared Ops/logging service"),
            ("Ex-ADL / Local HUB", "Scope add", "Application logs newly into Dynatrace"),
        ],
        col_widths=[Inches(4.0), Inches(2.2), Inches(6.3)],
    )
    add_textbox(
        s,
        Inches(0.4),
        Inches(5.3),
        Inches(12.5),
        Inches(1.2),
        "Legend (ARB): Green = Reuse Shared Service by AXA-GO · Other legend colors = Modification / New AXA / New external service.\nReplace this slide with the official OneAXA domain map graphic from the master template for the live ARB.",
        size=12,
        color=GRAY,
    )
    footer(s, 7)

    # ----- 8 D-PRA -----
    s = blank_slide(prs)
    section_header_bar(s, "Alignment with Digital Platform Reference Architecture")
    key_message_box(s, ["Using Dynatrace (SaaS) for application logging system."])
    add_table(
        s,
        Inches(0.4),
        Inches(2.2),
        Inches(12.5),
        [
            ("D-PRA layer", "Highlight on slide", "What it means for this project"),
            ("UI", "Partner / 3rd Party / SaaS", "Dynatrace is consumed as external SaaS"),
            ("Platforms / Cloud Infrastructure", "IaaS / PaaS / XaaS", "Cloud binding / Cloud Permit relevant"),
            ("Platforms / Operations", "Monitoring-Alerting / Capacity / BCP", "Dynatrace owns ops observability capability"),
            ("API product layer", "N/A on summary", "No new business API; logging platform change"),
        ],
        col_widths=[Inches(3.5), Inches(3.8), Inches(5.2)],
    )
    add_textbox(
        s,
        Inches(0.4),
        Inches(5.2),
        Inches(12.5),
        Inches(1.2),
        "Paste the official D-PRA wallpaper from the AXA Japan ARB master behind this table for the formal meeting.",
        size=12,
        color=GRAY,
    )
    footer(s, 8)

    # ----- 9 Data Architecture -----
    s = blank_slide(prs)
    section_header_bar(s, "Data Architecture")
    key_message_box(
        s,
        [
            "Dynatrace implements 3 layers: Ingestion, Store, Query processing.",
            "Log retention default = 3 months; adjustable per application (example MyAXA = 1 year).",
        ],
    )
    y = Inches(2.4)
    flow_box(s, Inches(0.5), y, Inches(2.6), Inches(1.4), "1 Application Logs", "Lambda · OpenPaaS\nMPI/POD · LocalHUB/Ex-ADL")
    arrow_right(s, Inches(3.2), y + Inches(0.55))
    flow_box(s, Inches(3.75), y, Inches(2.6), Inches(1.4), "2 Dynatrace Ingestion", "OneAXA / SaaS ingest path")
    arrow_right(s, Inches(6.45), y + Inches(0.55))
    flow_box(s, Inches(7.0), y, Inches(2.6), Inches(1.4), "3 Dynatrace Store", "Buckets / tables / views\ntimestamp + raw logs")
    arrow_right(s, Inches(9.7), y + Inches(0.55))
    flow_box(s, Inches(10.25), y, Inches(2.6), Inches(1.4), "4 Query Engine", "Search / analytics\nin Dynatrace")

    add_table(
        s,
        Inches(0.5),
        Inches(4.2),
        Inches(12.3),
        [
            ("Retention", "Value"),
            ("Default", "3 months"),
            ("Example longer retention", "MyAXA = 1 year"),
            ("Policy", "Adjust per application as needed"),
        ],
        col_widths=[Inches(4.0), Inches(8.3)],
    )
    footer(s, 9)

    # ----- 10 Security -----
    s = blank_slide(prs)
    section_header_bar(s, "Security Architecture")
    key_message_box(s, ["Using OneAccount login to Dynatrace (SaaS).", "IAM CoE Demand: Already integrated as provided by AXA-GO."])

    flow_box(s, Inches(1.0), Inches(2.6), Inches(3.2), Inches(1.5), "I&T members", "AXA-HQ Shirokane", fill=LIGHT_GRAY)
    arrow_right(s, Inches(4.4), Inches(3.2))
    flow_box(s, Inches(5.0), Inches(2.6), Inches(3.2), Inches(1.5), "AXA-GO OneAccount", "SSO / enterprise identity", fill=LIGHT_BLUE)
    arrow_right(s, Inches(8.4), Inches(3.2))
    flow_box(s, Inches(9.0), Inches(2.6), Inches(3.5), Inches(1.5), "Dynatrace Web Console", "RBAC enforced", fill=WHITE)

    add_textbox(s, Inches(9.5), Inches(4.4), Inches(3), Inches(0.3), "↓", size=18, bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)
    flow_box(s, Inches(9.0), Inches(4.8), Inches(3.5), Inches(1.2), "Dynatrace Store", "Log data under SaaS controls")

    add_table(
        s,
        Inches(1.0),
        Inches(4.8),
        Inches(7.5),
        [
            ("Control", "What it means"),
            ("OneAccount", "Enterprise SSO into Dynatrace"),
            ("RBAC", "Role-based access on console"),
            ("AXA-GO provided", "IAM integration not invented by this project"),
        ],
        col_widths=[Inches(2.5), Inches(5.0)],
    )
    footer(s, 10)

    # ----- 11 HLD overview -----
    s = blank_slide(prs)
    section_header_bar(s, "Technical Infrastructure — HLD Overview")
    key_message_box(
        s,
        [
            "HQ users reach Dynatrace Console/API/Grail over HTTPS via proxy.",
            "ActiveGates bridge Tokyo/Singapore MPI, OpenPaaS, LocalHUB, and AWS Firehose paths.",
        ],
    )
    add_table(
        s,
        Inches(0.4),
        Inches(2.2),
        Inches(12.5),
        [
            ("Zone", "Components", "Path"),
            ("AXA-HQ Shirokane", "IT members, Proxy", "HTTPS → Dynatrace Console / API / Grail"),
            ("AXA-GO Dynatrace", "SaaS tenant, Grail", "Central observability platform"),
            ("ActiveGate", "Regional ActiveGates", "TCP/HTTPS from hubs / MPI / apps"),
            ("OpenPaaS", "Multi Log Forwarder", "Egress → Dynatrace"),
            ("AWS Shared / MPI", "Firehose, CloudWatch, EC2 apps", "Cloud + MPI ingest paths"),
            ("Local HUBs", "TRS / ALJAGU / Ex-ADL", "App logs via ActiveGate path"),
            ("Support tools", "OneAccount, Config Repository", "VPN / HTTPS for employees"),
        ],
        col_widths=[Inches(2.8), Inches(4.2), Inches(5.5)],
    )
    add_textbox(
        s,
        Inches(0.4),
        Inches(6.3),
        Inches(12.5),
        Inches(0.5),
        "Replace with the detailed Visio/HLD graphic from the formal ARB pack for pixel-accurate ports and trust boundaries.",
        size=11,
        color=GRAY,
    )
    footer(s, 11)

    # ----- 12 Next / open items -----
    s = blank_slide(prs)
    section_header_bar(s, "Open Items and Next Steps")
    add_table(
        s,
        Inches(0.4),
        Inches(1.2),
        Inches(12.5),
        [
            ("Item", "Status on captured slides", "Action"),
            ("Security Review", "TBA – before Friday", "Close before ARB / update summary"),
            ("Architect Name on summary", "Blank", "Fill architect-in-charge name"),
            ("IR / CP dated approvals", "Placeholders", "Update after IR / Cloud Permit"),
            ("Impacted Applications list", "Mandatory; not in photo set", "Attach full app inventory slide"),
            ("Decommission Plan", "Mandatory for BP/Infra A", "Attach Splunk MPI decommission slide"),
            ("Cloud Binding Checklist", "Mandatory for CP", "Complete CP checklist slide"),
            ("Production", "Q3-2026", "Keep Gate2 cost 25M Yen aligned"),
        ],
        col_widths=[Inches(3.5), Inches(4.0), Inches(5.0)],
    )
    add_textbox(
        s,
        Inches(0.4),
        Inches(5.5),
        Inches(12.5),
        Inches(1.2),
        "This deck is an ARB-style working copy generated from screenshot capture of P260120F.\nFor the formal board, merge branding/master slides from the official AXA Japan ARB PowerPoint template (ver.5.2).",
        size=12,
        color=GRAY,
    )
    footer(s, 12)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    base = Path(__file__).resolve().parent
    out = base / "11-P260120F-Splunk-to-Dynatrace-ARB-style.pptx"
    build(out)
