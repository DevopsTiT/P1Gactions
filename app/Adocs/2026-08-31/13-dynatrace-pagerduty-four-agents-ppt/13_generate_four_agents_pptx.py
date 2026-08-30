#!/usr/bin/env python3
"""ARB-style PPT: Dynatrace → PagerDuty Four AI Agents (SRE/Scribe/Shift/Insights)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x2C, 0x5C)
NAVY_MID = RGBColor(0x14, 0x3D, 0x7A)
BLUE_HDR = RGBColor(0x1F, 0x4E, 0x8C)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
GREEN = RGBColor(0x1E, 0x84, 0x4E)
LIGHT_BLUE = RGBColor(0xD6, 0xE6, 0xF5)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x55, 0x55, 0x55)
YELLOW = RGBColor(0xFF, 0xF2, 0xCC)
PURPLE = RGBColor(0x5B, 0x2C, 0x6F)


def set_run(run, text, size=12, bold=False, color=BLACK, font_name="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def add_textbox(slide, left, top, width, height, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color)
    return box


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def add_round_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def footer(slide, page):
    add_textbox(
        slide,
        Inches(0.3),
        Inches(7.15),
        Inches(10),
        Inches(0.3),
        "ARB-style · Dynatrace → PagerDuty Four AI Agents",
        size=9,
        color=GRAY,
    )
    add_textbox(
        slide,
        Inches(11.5),
        Inches(7.15),
        Inches(1.5),
        Inches(0.3),
        f"{page}",
        size=10,
        bold=True,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
    )


def section_header_bar(slide, title):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), NAVY)
    add_textbox(slide, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.5), title, size=22, bold=True, color=WHITE)


def key_message_box(slide, lines, top=Inches(1.0)):
    h = Inches(0.55 + 0.24 * len(lines))
    add_rect(slide, Inches(0.35), top, Inches(12.6), h, LIGHT_BLUE, NAVY_MID)
    add_textbox(slide, Inches(0.5), top + Inches(0.06), Inches(1.5), Inches(0.28), "Key Message", size=11, bold=True, color=NAVY)
    y = top + Inches(0.32)
    for line in lines:
        add_textbox(slide, Inches(0.5), y, Inches(12.2), Inches(0.26), f"• {line}", size=12, color=BLACK)
        y += Inches(0.24)
    return y + Inches(0.1)


def add_table(slide, left, top, width, rows, col_widths=None, header=True):
    cols = len(rows[0])
    table_shape = slide.shapes.add_table(len(rows), cols, left, top, width, Inches(0.34 * len(rows)))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            is_hdr = header and r_idx == 0
            set_run(run, str(cell_text), size=10 if not is_hdr else 11, bold=is_hdr, color=WHITE if is_hdr else BLACK)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE_HDR if is_hdr else (LIGHT_GRAY if r_idx % 2 else WHITE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape


def flow_box(slide, left, top, width, height, title, body, fill=WHITE, title_color=NAVY):
    shape = add_round_rect(slide, left, top, width, height, fill, NAVY_MID)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    set_run(r0, title, size=11, bold=True, color=title_color)
    p1 = tf.add_paragraph()
    r1 = p1.add_run()
    set_run(r1, body, size=9, color=GRAY)
    return shape


def arrow_right(slide, left, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, Inches(0.42), Inches(0.22))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY_MID
    shape.line.fill.background()
    return shape


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build(out_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1 Title
    s = blank_slide(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_textbox(s, Inches(0.6), Inches(0.45), Inches(8), Inches(0.35), "Architecture Design Proposal", size=14, bold=True, color=LIGHT_BLUE)
    add_textbox(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.9), "Dynatrace → PagerDuty", size=34, bold=True, color=WHITE)
    add_textbox(s, Inches(0.6), Inches(2.05), Inches(12), Inches(0.5), "Four AI Agents Architecture", size=26, bold=True, color=LIGHT_BLUE)
    add_textbox(
        s,
        Inches(0.6),
        Inches(2.7),
        Inches(12),
        Inches(0.4),
        "SRE Agent · Scribe Agent · Shift Agent · Insights Agent",
        size=16,
        color=WHITE,
    )
    meta = [
        ("Field", "Value"),
        ("Objective", "Architecture + POC design for PagerDuty Advance agents fed by Dynatrace"),
        ("Platform", "Dynatrace (detect) → PagerDuty (page) → Advance AI Agents (assist)"),
        ("Umbrella product", "PagerDuty Advance (AI Actions)"),
        ("Chat surface", "Microsoft Teams (primary) and/or Slack; PD web for SRE tab"),
        ("Safe demo rule", "Test Service + Level-1 test schedule — never prod pages"),
        ("Related", "Builds on four-agents POC + Teams-only enablement notes"),
    ]
    add_table(s, Inches(0.6), Inches(3.3), Inches(11.8), meta, col_widths=[Inches(2.8), Inches(9.0)])
    add_textbox(s, Inches(0.6), Inches(6.85), Inches(11), Inches(0.3), "ARB-style working deck · LAST UPDATED 2026/08/31", size=10, color=LIGHT_BLUE)

    # 2 Summary
    s = blank_slide(prs)
    section_header_bar(s, "Architecture Design <Summary Sheet>")
    add_table(
        s,
        Inches(0.35),
        Inches(1.1),
        Inches(12.6),
        [
            ("Item", "Answer"),
            ("Objective", "Prove Dynatrace problems become clean PD incidents, then four Advance agents help humans"),
            ("Background WHY", "Dynatrace finds the problem; humans still dig docs, take notes, fix coverage, and report trends. Agents cut that toil."),
            ("Design Overview", "Dynatrace Problem → PD Service (optional Event Orchestration) → page on-call → SRE triage / Scribe bridge / Shift coverage / Insights analytics"),
            ("Not one of the four", "Advance Assistant (router), Event Intelligence, PIR drafts — related, not the named suite"),
            ("Cloud / SaaS", "Dynatrace SaaS + PagerDuty Advance"),
            ("Chat", "Teams and/or Slack; Teams-only: SRE/Scribe/Insights Q&A full; Shift DMs Slack-first"),
            ("Cost model", "SRE 4 Actions/ask; Scribe ~6/30min +2 summary; Shift 0; Insights 0"),
            ("Safety", "Test Service / test schedule only; human confirms remediations"),
        ],
        col_widths=[Inches(2.8), Inches(9.8)],
    )
    footer(s, 2)

    # 3 Decision / which agent
    s = blank_slide(prs)
    section_header_bar(s, "Which Agent Do You Need?")
    key_message_box(s, ["Pick the agent from the job — not from the logo."])
    add_table(
        s,
        Inches(0.35),
        Inches(2.1),
        Inches(12.6),
        [
            ("Job", "Agent", "When"),
            ("Active incident triage / RCA / runbook", "SRE Agent", "First 10 minutes after Dynatrace pages"),
            ("Bridge notes / PIR draft", "Scribe Agent", "War room Zoom/Teams/Meet"),
            ("OOO vs on-call / coverage override", "Shift Agent", "Before the next page; Level-1 schedules"),
            ("MTTR / MTTA / volume trends", "Insights Agent", "After incidents; leadership prep"),
        ],
        col_widths=[Inches(5.0), Inches(2.8), Inches(4.8)],
    )
    add_textbox(
        s,
        Inches(0.4),
        Inches(5.3),
        Inches(12.5),
        Inches(1.2),
        "Prerequisite gate: PagerDuty Advance enabled → connect Teams/Slack → enable agents → TEST service only.\nWithout Advance, agent toggles do nothing.",
        size=13,
        color=GRAY,
    )
    footer(s, 3)

    # 4 AS-IS
    s = blank_slide(prs)
    section_header_bar(s, "Solution Context — AS-IS")
    key_message_box(
        s,
        [
            "Dynatrace detects and creates/notifies a PagerDuty incident.",
            "Humans do triage, notes, coverage, and analytics mostly by hand.",
        ],
    )
    add_round_rect(s, Inches(0.4), Inches(2.25), Inches(1.5), Inches(0.5), ORANGE)
    add_textbox(s, Inches(0.5), Inches(2.32), Inches(1.3), Inches(0.35), "AS-IS", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    flow_box(s, Inches(0.4), Inches(3.1), Inches(2.5), Inches(1.4), "Dynatrace", "Problem / Davis AI\ncreates signal")
    arrow_right(s, Inches(3.0), Inches(3.7))
    flow_box(s, Inches(3.5), Inches(3.1), Inches(2.6), Inches(1.4), "PagerDuty", "Incident + escalation\npages on-call")
    arrow_right(s, Inches(6.2), Inches(3.7))
    flow_box(s, Inches(6.7), Inches(3.1), Inches(2.8), Inches(1.4), "Human on-call", "Dig docs · type notes\nAsk coverage · Export CSV", fill=YELLOW)
    flow_box(s, Inches(9.9), Inches(3.1), Inches(2.9), Inches(1.4), "Pain", "Slow MTTA/MTTR\nNoise · missed notes\nSpreadsheet coverage")
    add_textbox(
        s,
        Inches(0.4),
        Inches(5.0),
        Inches(12.5),
        Inches(1.5),
        "AS-IS gap: Dynatrace is strong at detect. PagerDuty is strong at page.\nThe middle toil (triage, bridge notes, coverage, trends) is still mostly human.",
        size=13,
        color=GRAY,
    )
    footer(s, 4)

    # 5 TO-BE
    s = blank_slide(prs)
    section_header_bar(s, "Solution Context — TO-BE")
    key_message_box(
        s,
        [
            "Same Dynatrace → PagerDuty path, plus four Advance agents beside the human.",
            "Agents suggest and assist; humans still own blast radius.",
        ],
    )
    add_round_rect(s, Inches(0.4), Inches(2.25), Inches(1.5), Inches(0.5), GREEN)
    add_textbox(s, Inches(0.5), Inches(2.32), Inches(1.3), Inches(0.35), "TO-BE", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    flow_box(s, Inches(0.3), Inches(3.0), Inches(2.2), Inches(1.2), "Dynatrace", "Problem URL\nin PD payload")
    arrow_right(s, Inches(2.55), Inches(3.5))
    flow_box(s, Inches(3.05), Inches(3.0), Inches(2.3), Inches(1.2), "PagerDuty", "Service + optional\nEvent Orchestration")
    arrow_right(s, Inches(5.45), Inches(3.5))
    hub = add_round_rect(s, Inches(5.95), Inches(2.85), Inches(2.7), Inches(1.5), WHITE, ACCENT_RED)
    tf = hub.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, "Advance Agents", size=14, bold=True, color=ACCENT_RED)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    set_run(r2, "SRE · Scribe\nShift · Insights", size=11, color=GRAY)
    arrow_right(s, Inches(8.8), Inches(3.5))
    flow_box(s, Inches(9.3), Inches(3.0), Inches(3.5), Inches(1.2), "Human decides", "Confirm remediation\nOwn the page")

    flow_box(s, Inches(0.3), Inches(4.6), Inches(2.9), Inches(1.5), "SRE", "Triage + runbook\n4 Actions / ask", fill=LIGHT_BLUE)
    flow_box(s, Inches(3.4), Inches(4.6), Inches(2.9), Inches(1.5), "Scribe", "Meeting transcript\n~6/30m + 2 summary", fill=LIGHT_BLUE)
    flow_box(s, Inches(6.5), Inches(4.6), Inches(2.9), Inches(1.5), "Shift", "Coverage / override\n0 Actions", fill=LIGHT_BLUE)
    flow_box(s, Inches(9.6), Inches(4.6), Inches(3.2), Inches(1.5), "Insights", "MTTR Q&A + tips\n0 Actions", fill=LIGHT_BLUE)
    footer(s, 5)

    # 6 E2E data flow
    s = blank_slide(prs)
    section_header_bar(s, "End-to-End Data Flow")
    key_message_box(s, ["Detect → Notify → Page → Assist → Resolve → Learn"])
    add_table(
        s,
        Inches(0.35),
        Inches(2.1),
        Inches(12.6),
        [
            ("Stage", "System", "Agent touchpoint"),
            ("1 Detect", "Dynatrace Davis Problem", "None yet — keep problem URL in the event"),
            ("2 Notify", "Dynatrace → PD Events / Workflow", "Clean routing key + urgency matter"),
            ("3 Optional noise cut", "PD Event Orchestration / grouping", "Fewer junk incidents before agents run"),
            ("4 Page", "Escalation Policy", "Shift kept coverage valid beforehand"),
            ("5 Triage", "Teams/Slack/PD web", "SRE Agent summarizes + runbook + next steps"),
            ("6 Collaborate", "Zoom / Teams / Meet bridge", "Scribe joins, transcript, summary, PIR context"),
            ("7 Resolve", "PD Resolve", "SRE service memory updates"),
            ("8 Improve", "Analytics / chat next week", "Insights MTTR/MTTA Q&A (weekly DMs Slack-first)"),
        ],
        col_widths=[Inches(2.4), Inches(4.2), Inches(6.0)],
    )
    footer(s, 6)

    # 7 Four agents matrix
    s = blank_slide(prs)
    section_header_bar(s, "Four Agents at a Glance")
    add_table(
        s,
        Inches(0.25),
        Inches(1.1),
        Inches(12.8),
        [
            ("Agent", "What it is", "Primary surface", "AI Actions", "Teams-only"),
            ("SRE", "Virtual responder: context, runbook, next steps", "Teams @pagerduty / PD SRE tab", "4 / ask", "Full (EA)"),
            ("Scribe", "Joins meeting; transcript + wrap-up", "Teams meeting + advance scribe", "~6/30m +2", "Strong"),
            ("Shift", "OOO conflict + coverage override", "Slack DMs / web+Calendar Path B", "0", "Partial"),
            ("Insights", "Analytics Q&A + weekly maturity tips", "Teams @pagerduty Q&A", "0", "Q&A yes; weekly DM Slack"),
        ],
        col_widths=[Inches(1.5), Inches(4.0), Inches(3.5), Inches(1.5), Inches(2.3)],
    )
    add_textbox(
        s,
        Inches(0.4),
        Inches(5.0),
        Inches(12.5),
        Inches(1.5),
        "Dynatrace tip: Dynatrace usually CREATES the incident (ingress).\nIt is not the main item on the SRE Agent connector list — enrich with Grafana/Datadog/CloudWatch/Confluence/GitHub as needed.\nKeep the Dynatrace problem URL in custom_details so humans and SRE Agent can deep-link.",
        size=12,
        color=GRAY,
    )
    footer(s, 7)

    # 8 SRE detail
    s = blank_slide(prs)
    section_header_bar(s, "1) SRE Agent — Live Triage")
    key_message_box(
        s,
        [
            "Reads incident + ~2k chars custom_details + runbook + optional connectors.",
            "Suggests next steps; human must confirm remediations. Memory saves on Resolve.",
        ],
    )
    add_table(
        s,
        Inches(0.35),
        Inches(2.2),
        Inches(12.6),
        [
            ("Surface", "How to start", "Good first ask"),
            ("MS Teams (Early Access)", "@pagerduty … in mapped channel", "What are some likely root causes?"),
            ("Incident SRE tab", "PD web → incident → SRE Agent", "Analyze past incidents"),
            ("Ops Console", "AIOps + Advance → SRE tab", "What steps should I take first?"),
            ("Virtual responder", "Incident Workflow / Escalation (EA)", "Same 4 Actions per trigger"),
        ],
        col_widths=[Inches(3.2), Inches(5.0), Inches(4.4)],
    )
    add_textbox(
        s,
        Inches(0.4),
        Inches(5.3),
        Inches(12.5),
        Inches(1.3),
        "POC story: Dynatrace “Checkout latency high” on test Service → Teams ask → upload poc-checkout-latency.md\n→ past incidents → first steps → READ any remediation → Resolve.\nLimits: runbook ≤100 KB; 25 files/conversation; no customer PII in channel questions.",
        size=12,
        color=GRAY,
    )
    footer(s, 8)

    # 9 Scribe
    s = blank_slide(prs)
    section_header_bar(s, "2) Scribe Agent — Bridge Notes")
    key_message_box(
        s,
        [
            "Joins Zoom / Teams / Meet; streams transcript; posts decisions / actions / attendees.",
            "Human must join within 15 minutes for auto-join. Cap: 1 Scribe/meeting; 10 concurrent.",
        ],
    )
    add_table(
        s,
        Inches(0.35),
        Inches(2.2),
        Inches(12.6),
        [
            ("Step", "Action"),
            ("1", "Enable Scribe in AI Settings (US often manual; EU may default on)"),
            ("2", "Put Teams/Zoom/Meet URL (with passcode) on the test incident"),
            ("3", "Auto-join OR Teams: @PagerDuty advance scribe"),
            ("4", "Join yourself within 15 minutes; admit lobby if needed"),
            ("5", "Speak symptom / cause / decision clearly"),
            ("6", "End meeting → confirm summary → optional PIR draft"),
        ],
        col_widths=[Inches(1.2), Inches(11.4)],
    )
    footer(s, 9)

    # 10 Shift + Insights
    s = blank_slide(prs)
    section_header_bar(s, "3) Shift + 4) Insights")
    add_textbox(s, Inches(0.4), Inches(1.05), Inches(6), Inches(0.3), "Shift Agent", size=16, bold=True, color=NAVY)
    add_table(
        s,
        Inches(0.35),
        Inches(1.4),
        Inches(12.6),
        [
            ("Topic", "Detail"),
            ("Job", "Detect OOO vs Level-1 on-call; request coverage; write override"),
            ("Cost", "0 AI Actions"),
            ("Slack path", "Conflict DM → Request coverage → accept → schedule update"),
            ("Teams-only", "Path B recommended: Google Calendar + PD web manual override"),
            ("Safety", "Test schedule / test EP only — never prod primary"),
        ],
        col_widths=[Inches(2.5), Inches(10.1)],
    )
    add_textbox(s, Inches(0.4), Inches(4.5), Inches(6), Inches(0.3), "Insights Agent", size=16, bold=True, color=NAVY)
    add_table(
        s,
        Inches(0.35),
        Inches(4.85),
        Inches(12.6),
        [
            ("Topic", "Detail"),
            ("Job", "Conversational MTTR/MTTA/volume; weekly maturity tips"),
            ("Cost", "0 AI Actions"),
            ("Teams", "On-demand @pagerduty Q&A is GA; sanity-check vs Analytics UI"),
            ("Slack-first gap", "Weekly proactive recommendation DMs still documented as Slack"),
        ],
        col_widths=[Inches(2.5), Inches(10.1)],
    )
    footer(s, 10)

    # 11 Enablement + cost
    s = blank_slide(prs)
    section_header_bar(s, "Enablement, Cost, Safety")
    key_message_box(s, ["Enable Advance once. Demo on a test Service that pages only you."])
    add_table(
        s,
        Inches(0.35),
        Inches(2.05),
        Inches(12.6),
        [
            ("Layer", "Do this"),
            ("0 Shared", "Advance On · Teams/Slack Connected · linkUser · test Service map"),
            ("Graph (Teams)", "ChatMessage.Read(.All) for Advance; prefer User.ReadBasic.All"),
            ("Dynatrace", "Test integration / routing key only; problem URL in payload"),
            ("SRE connectors", "Optional one real connector (not Dynatrace as “connector list” item)"),
            ("Cost", "SRE 4/ask · Scribe ~6/30m+2 · Shift 0 · Insights 0"),
            ("Safety", "No prod escalation · confirm remediations · no PII in channel asks"),
        ],
        col_widths=[Inches(2.5), Inches(10.1)],
    )
    footer(s, 11)

    # 12 POC plan
    s = blank_slide(prs)
    section_header_bar(s, "Suggested POC Afternoon + Open Items")
    add_table(
        s,
        Inches(0.35),
        Inches(1.1),
        Inches(12.6),
        [
            ("Order", "Block", "Time box"),
            ("0", "Shared enablement (Admins)", "30–60 min"),
            ("1", "SRE Agent on Dynatrace-shaped test incident", "15–20 min"),
            ("2", "Scribe short Teams meeting (≤10 min)", "15–20 min"),
            ("3", "Shift Path B (Calendar + web override)", "10–15 min"),
            ("4", "Insights conversational Q&A", "10–15 min"),
        ],
        col_widths=[Inches(1.3), Inches(7.5), Inches(3.8)],
    )
    add_table(
        s,
        Inches(0.35),
        Inches(4.3),
        Inches(12.6),
        [
            ("Open item", "Note"),
            ("Advance entitlement", "Need trial/add-on before toggles work"),
            ("Teams Early Access for SRE", "Confirm tenant has Teams EA"),
            ("Shift full DM path", "Needs Slack — document if Teams-only"),
            ("Insights weekly DMs", "Slack-first — document gap"),
            ("Prod keys", "Never point Dynatrace prod at POC Services"),
        ],
        col_widths=[Inches(3.5), Inches(9.1)],
    )
    footer(s, 12)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    base = Path(__file__).resolve().parent
    build(base / "13-Dynatrace-PagerDuty-Four-AI-Agents-ARB-style.pptx")
