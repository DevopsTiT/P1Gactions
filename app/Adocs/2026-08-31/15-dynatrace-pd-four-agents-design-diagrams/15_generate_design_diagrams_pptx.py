#!/usr/bin/env python3
"""ARB-style Solution Context / Impact / Alignment / HLD diagrams
for Dynatrace → PagerDuty Four AI Agents (design pack)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x2C, 0x5C)
NAVY_MID = RGBColor(0x14, 0x3D, 0x7A)
BLUE_HDR = RGBColor(0x1F, 0x4E, 0x8C)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
GREEN = RGBColor(0x1E, 0x84, 0x4E)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
LIGHT_BLUE = RGBColor(0xD6, 0xE6, 0xF5)
LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
LIGHT_ORANGE = RGBColor(0xFD, 0xE8, 0xD0)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF7)
LIGHT_RED = RGBColor(0xFA, 0xE5, 0xD3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x55, 0x55, 0x55)
YELLOW = RGBColor(0xFF, 0xF2, 0xCC)


def set_run(run, text, size=12, bold=False, color=BLACK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_textbox(slide, left, top, width, height, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color)
    return box


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def card(slide, left, top, width, height, title, body, fill=WHITE, title_color=NAVY, line=NAVY_MID):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.08)
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    set_run(r0, title, size=12, bold=True, color=title_color)
    if body:
        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        set_run(r1, body, size=10, color=GRAY)
    return shape


def arrow(slide, left, top, width=Inches(0.45)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.22))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY_MID
    shape.line.fill.background()
    return shape


def down_arrow(slide, left, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, Inches(0.28), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY_MID
    shape.line.fill.background()
    return shape


def header(slide, title):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.8), NAVY)
    add_textbox(slide, Inches(0.35), Inches(0.18), Inches(12.5), Inches(0.45), title, size=22, bold=True, color=WHITE)


def key_msg(slide, lines):
    add_rect(slide, Inches(0.3), Inches(0.95), Inches(12.7), Inches(0.55 + 0.22 * (len(lines) - 1)), LIGHT_BLUE, NAVY_MID)
    add_textbox(slide, Inches(0.45), Inches(1.0), Inches(1.5), Inches(0.25), "Key Message", size=11, bold=True, color=NAVY)
    y = Inches(1.22)
    for line in lines:
        add_textbox(slide, Inches(0.45), y, Inches(12.4), Inches(0.22), f"• {line}", size=12, color=BLACK)
        y += Inches(0.22)


def footer(slide, page):
    add_textbox(slide, Inches(0.3), Inches(7.15), Inches(10), Inches(0.25), "AXA Japan ARB Template style · Dynatrace → PagerDuty Four AI Agents · Design", size=9, color=GRAY)
    add_textbox(slide, Inches(11.6), Inches(7.15), Inches(1.4), Inches(0.25), str(page), size=10, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def badge(slide, left, top, text, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.5), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text, size=14, bold=True, color=WHITE)
    return shape


def build(out: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ----- 1 Title -----
    s = blank(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_textbox(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.4), "Architecture Design · Solution Context Pack", size=14, bold=True, color=LIGHT_BLUE)
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(12), Inches(0.8), "Dynatrace → PagerDuty", size=36, bold=True, color=WHITE)
    add_textbox(s, Inches(0.5), Inches(2.3), Inches(12), Inches(0.5), "Four AI Agents — Design Diagrams", size=24, bold=True, color=LIGHT_BLUE)
    add_textbox(s, Inches(0.5), Inches(3.1), Inches(12), Inches(0.4), "SRE · Scribe · Shift · Insights  |  ARB-style AS-IS / TO-BE / Impact / HLD", size=16, color=WHITE)
    card(s, Inches(0.5), Inches(4.0), Inches(5.8), Inches(2.0), "Purpose", "Same diagram style as Splunk→Dynatrace\nSolution Context slides — for design review\nand POC architecture discussion.", fill=RGBColor(0x12, 0x3A, 0x6B), title_color=WHITE, line=LIGHT_BLUE)
    card(s, Inches(6.6), Inches(4.0), Inches(5.8), Inches(2.0), "Contents", "1 Title\n2 AS-IS Solution Context\n3 TO-BE Solution Context\n4 Impacted Platforms\n5 Alignment (Ops / IR)\n6 Technical HLD\n7 Agent hub detail\n8 Design checklist", fill=RGBColor(0x12, 0x3A, 0x6B), title_color=WHITE, line=LIGHT_BLUE)

    # ----- 2 AS-IS -----
    s = blank(prs)
    header(s, "Solution Context Diagram — AS-IS")
    key_msg(
        s,
        [
            "Dynatrace detects problems and notifies PagerDuty; escalation pages on-call.",
            "Triage, bridge notes, coverage, and analytics are still mostly manual human work.",
        ],
    )
    badge(s, Inches(0.35), Inches(1.85), "AS-IS", ORANGE)

    # Top row sources
    card(s, Inches(0.4), Inches(2.5), Inches(2.6), Inches(1.25), "Dynatrace", "Davis Problem\nmetrics / logs / traces", fill=LIGHT_ORANGE, title_color=ORANGE)
    arrow(s, Inches(3.1), Inches(3.0))
    card(s, Inches(3.65), Inches(2.5), Inches(2.7), Inches(1.25), "PagerDuty", "Incident created\nEscalation Policy", fill=LIGHT_BLUE)
    arrow(s, Inches(6.45), Inches(3.0))
    card(s, Inches(7.0), Inches(2.4), Inches(2.9), Inches(1.45), "Human on-call", "Ack · dig docs · type notes\nAsk coverage · export CSV", fill=YELLOW, title_color=ORANGE, line=ORANGE)

    # Pain boxes
    card(s, Inches(10.2), Inches(2.4), Inches(2.8), Inches(1.45), "Pain", "Slow first 10 min\nLost bridge decisions\nCoverage gaps\nManual trends", fill=LIGHT_RED, title_color=ACCENT_RED, line=ACCENT_RED)

    # Bottom supporting
    card(s, Inches(0.4), Inches(4.2), Inches(2.6), Inches(1.2), "Apps / Infra", "Checkout · hosts · DB\n(monitored by Dynatrace)", fill=WHITE)
    card(s, Inches(3.3), Inches(4.2), Inches(2.8), Inches(1.2), "Chat (optional)", "Teams / Slack cards only\nNo AI agents yet", fill=WHITE)
    card(s, Inches(6.4), Inches(4.2), Inches(3.0), Inches(1.2), "Schedule", "Level-1 on-call\nManual OOO overrides", fill=WHITE)
    card(s, Inches(9.7), Inches(4.2), Inches(3.2), Inches(1.2), "Analytics", "PD Analytics UI\nManual pull for leadership", fill=WHITE)

    add_textbox(s, Inches(0.4), Inches(5.7), Inches(12.5), Inches(0.9), "AS-IS gap: Detect (Dynatrace) and Page (PagerDuty) exist. The assist layer between page and resolve is human-only.", size=13, bold=True, color=NAVY)
    footer(s, 2)

    # ----- 3 TO-BE -----
    s = blank(prs)
    header(s, "Solution Context Diagram — TO-BE")
    key_msg(
        s,
        [
            "PagerDuty Advance Four AI Agents assist beside the human after Dynatrace creates the incident.",
            "Optional Event Orchestration cuts noise before agents run. Humans still confirm remediations.",
        ],
    )
    badge(s, Inches(0.35), Inches(1.85), "TO-BE", GREEN)

    # Sources into hub
    card(s, Inches(0.3), Inches(2.45), Inches(2.4), Inches(1.1), "Dynatrace", "Problem + problem URL\nin PD payload", fill=LIGHT_GREEN, title_color=GREEN)
    card(s, Inches(0.3), Inches(3.7), Inches(2.4), Inches(1.0), "Optional EO", "Route / pause / group\n(noise reduction)", fill=LIGHT_GRAY)
    card(s, Inches(0.3), Inches(4.9), Inches(2.4), Inches(1.0), "Ops (HQ)", "Teams / PD web\nOneAccount-style access", fill=WHITE)

    arrow(s, Inches(2.8), Inches(2.9))
    arrow(s, Inches(2.8), Inches(4.05))
    arrow(s, Inches(2.8), Inches(5.25))

    # Center hub - PagerDuty + Advance
    hub = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.4), Inches(2.7), Inches(3.3), Inches(2.6))
    hub.fill.solid()
    hub.fill.fore_color.rgb = WHITE
    hub.line.color.rgb = ACCENT_RED
    hub.line.width = Pt(2.5)
    tf = hub.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, "PagerDuty", size=14, bold=True, color=ACCENT_RED)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    set_run(r2, "Incident hub", size=11, color=GRAY)
    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    set_run(r3, "\nAdvance AI Agents", size=13, bold=True, color=NAVY)
    p4 = tf.add_paragraph()
    r4 = p4.add_run()
    set_run(r4, "SRE · Scribe\nShift · Insights", size=12, bold=True, color=TEAL)

    arrow(s, Inches(6.85), Inches(3.8))

    # Right side agents detail
    card(s, Inches(7.4), Inches(2.4), Inches(2.7), Inches(1.05), "SRE Agent", "Triage + runbook\n4 Actions / ask", fill=LIGHT_BLUE)
    card(s, Inches(10.3), Inches(2.4), Inches(2.7), Inches(1.05), "Scribe Agent", "Meeting transcript\n~6/30m + 2", fill=LIGHT_BLUE)
    card(s, Inches(7.4), Inches(3.6), Inches(2.7), Inches(1.05), "Shift Agent", "Coverage / override\n0 Actions", fill=LIGHT_BLUE)
    card(s, Inches(10.3), Inches(3.6), Inches(2.7), Inches(1.05), "Insights Agent", "MTTR Q&A\n0 Actions", fill=LIGHT_BLUE)

    card(s, Inches(7.4), Inches(4.9), Inches(5.6), Inches(1.15), "Human decides", "Confirm remediations · Own blast radius · Resolve → SRE memory", fill=YELLOW, title_color=ORANGE, line=ORANGE)
    footer(s, 3)

    # ----- 4 Impacted Platforms -----
    s = blank(prs)
    header(s, "Impacted Platforms")
    key_msg(s, ["Four AI Agents consolidate assist work on the incident path and speed troubleshooting after Dynatrace pages."])

    # Legend
    add_textbox(s, Inches(0.4), Inches(1.85), Inches(12), Inches(0.3), "Legend (ARB-style)", size=12, bold=True, color=NAVY)
    card(s, Inches(0.4), Inches(2.2), Inches(2.8), Inches(0.7), "Green", "Reuse existing\nDynatrace / PD / chat", fill=LIGHT_GREEN, title_color=GREEN, line=GREEN)
    card(s, Inches(3.4), Inches(2.2), Inches(2.8), Inches(0.7), "Orange", "Modify / enable\nAdvance + agents", fill=LIGHT_ORANGE, title_color=ORANGE, line=ORANGE)
    card(s, Inches(6.4), Inches(2.2), Inches(3.0), Inches(0.7), "Blue", "Primary design\nsurface", fill=LIGHT_BLUE, title_color=NAVY, line=NAVY)
    card(s, Inches(9.6), Inches(2.2), Inches(3.2), Inches(0.7), "Gray", "Out of scope\napp business logic", fill=LIGHT_GRAY, title_color=GRAY, line=GRAY)

    # Platform map simplified
    y = Inches(3.2)
    card(s, Inches(0.4), y, Inches(3.0), Inches(1.35), "Dynatrace", "Reuse — detect ingress\nKeep problem URL in event\nTest key for POC", fill=LIGHT_GREEN, title_color=GREEN, line=GREEN)
    card(s, Inches(3.6), y, Inches(3.0), Inches(1.35), "PagerDuty Core", "Modify — test Service/EP\nLevel-1 test schedule\nIntegration key hygiene", fill=LIGHT_ORANGE, title_color=ORANGE, line=ORANGE)
    card(s, Inches(6.8), y, Inches(3.0), Inches(1.35), "PagerDuty Advance", "Primary — enable agents\nAI Actions budget\nTeam Advance access", fill=LIGHT_BLUE, title_color=NAVY, line=NAVY)
    card(s, Inches(10.0), y, Inches(2.9), Inches(1.35), "Teams / Slack", "Modify — bot + linkUser\nGraph consent (Teams)\nStandard channel map", fill=LIGHT_ORANGE, title_color=ORANGE, line=ORANGE)

    y2 = Inches(4.75)
    card(s, Inches(0.4), y2, Inches(3.0), Inches(1.2), "IT / IT Platform", "Highlighted impact\nIR tooling change", fill=LIGHT_BLUE, title_color=NAVY, line=ACCENT_RED)
    card(s, Inches(3.6), y2, Inches(3.0), Inches(1.2), "Cloud / SaaS Ops", "Monitoring-Alerting layer\nDynatrace + PD SaaS", fill=LIGHT_BLUE, title_color=NAVY, line=ACCENT_RED)
    card(s, Inches(6.8), y2, Inches(3.0), Inches(1.2), "Calendar (optional)", "Google Calendar Ext.\nfor Shift conflicts", fill=LIGHT_GREEN, title_color=GREEN, line=GREEN)
    card(s, Inches(10.0), y2, Inches(2.9), Inches(1.2), "Business Apps", "No code rewrite\nfor product agents", fill=LIGHT_GRAY, title_color=GRAY, line=GRAY)
    footer(s, 4)

    # ----- 5 Alignment -----
    s = blank(prs)
    header(s, "Alignment with Ops / IR Reference Architecture")
    key_msg(s, ["Using Dynatrace (detect) + PagerDuty Advance agents (assist) on the incident response path."])

    # Layered stack like D-PRA ops highlight
    layers = [
        (Inches(2.3), "Channel / Users", "On-call · Commander · Leadership", LIGHT_GRAY),
        (Inches(3.15), "UI / Chat", "Microsoft Teams · Slack · PD web · Ops Console", LIGHT_BLUE),
        (Inches(4.0), "Assist (NEW)", "SRE · Scribe · Shift · Insights  (PagerDuty Advance)", YELLOW),
        (Inches(4.85), "Page / Incident", "PagerDuty Service · Escalation · Schedule · Incident", LIGHT_ORANGE),
        (Inches(5.7), "Detect", "Dynatrace Problems · optional Event Orchestration", LIGHT_GREEN),
    ]
    add_textbox(s, Inches(0.5), Inches(1.85), Inches(12), Inches(0.3), "Reference stack (design view)", size=13, bold=True, color=NAVY)
    for top, title, body, fill in layers:
        card(s, Inches(1.5), top, Inches(10.3), Inches(0.75), title, body, fill=fill, line=NAVY_MID)

    add_textbox(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4), "D-PRA analogy: Partner/SaaS + Ops Monitoring-Alerting layer — agents sit in Assist, not in app API product layer.", size=12, color=GRAY)
    footer(s, 5)

    # ----- 6 HLD -----
    s = blank(prs)
    header(s, "Technical Infrastructure — HLD Overview")
    key_msg(
        s,
        [
            "Dynatrace sends problems to PagerDuty; Advance agents attach to the incident and chat/meeting surfaces.",
            "Teams path needs Graph message-read for Advance; Shift full DMs prefer Slack.",
        ],
    )

    # Left users
    card(s, Inches(0.3), Inches(2.2), Inches(2.5), Inches(1.3), "Responder", "Teams / Slack\nlinkUser · graphAuth*", fill=WHITE)
    card(s, Inches(0.3), Inches(3.7), Inches(2.5), Inches(1.3), "Admin", "PD AI Settings\nMS Graph consent", fill=WHITE)
    card(s, Inches(0.3), Inches(5.2), Inches(2.5), Inches(1.1), "Calendar", "Google Calendar Ext.\n(Shift)", fill=LIGHT_GRAY)

    arrow(s, Inches(2.9), Inches(2.75))
    arrow(s, Inches(2.9), Inches(4.25))

    # Center PD
    card(s, Inches(3.5), Inches(2.2), Inches(3.2), Inches(2.0), "PagerDuty", "Service · Incident\nAdvance Agents\nAnalytics", fill=LIGHT_BLUE, title_color=NAVY, line=ACCENT_RED)
    card(s, Inches(3.5), Inches(4.4), Inches(3.2), Inches(1.5), "Optional EO", "Event Orchestration\nAlert Grouping", fill=LIGHT_GRAY)

    arrow(s, Inches(6.85), Inches(3.0))

    # Dynatrace
    card(s, Inches(7.4), Inches(2.2), Inches(2.7), Inches(1.5), "Dynatrace", "Problem / Davis\nNotify → PD key\nproblem URL", fill=LIGHT_GREEN, title_color=GREEN, line=GREEN)
    card(s, Inches(7.4), Inches(3.9), Inches(2.7), Inches(1.3), "Apps / Infra", "Checkout · hosts\nDB · JVM (monitored)", fill=WHITE)

    # Right surfaces
    card(s, Inches(10.3), Inches(2.2), Inches(2.7), Inches(1.2), "Teams meeting", "Scribe join\ntranscript", fill=LIGHT_BLUE)
    card(s, Inches(10.3), Inches(3.55), Inches(2.7), Inches(1.2), "Teams / Slack chat", "@pagerduty asks\nSRE · Insights", fill=LIGHT_BLUE)
    card(s, Inches(10.3), Inches(4.9), Inches(2.7), Inches(1.2), "PD web", "SRE Agent tab\nSchedules override", fill=LIGHT_ORANGE, title_color=ORANGE, line=ORANGE)

    add_textbox(s, Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.4), "* graphAuth only if tenant is Delegated-heavy. Prefer Application ChatMessage.Read.All for Advance.", size=11, color=GRAY)
    footer(s, 6)

    # ----- 7 Agent hub detail -----
    s = blank(prs)
    header(s, "Agent Hub Detail — Four Agents on One Incident")
    key_msg(s, ["One Dynatrace-fed incident can use multiple agents; pick by job."])

    # Center incident
    card(s, Inches(4.9), Inches(2.5), Inches(3.5), Inches(1.6), "PD Incident", "from Dynatrace Problem\n+ custom_details / problem URL", fill=YELLOW, title_color=ORANGE, line=ORANGE)

    card(s, Inches(0.4), Inches(2.3), Inches(3.8), Inches(1.9), "SRE Agent", "Summarize · runbook · past incidents\nNext steps · human confirm\nCost: 4 Actions / ask\nTeams EA + PD web tab", fill=LIGHT_BLUE)
    card(s, Inches(8.9), Inches(2.3), Inches(4.0), Inches(1.9), "Scribe Agent", "Join Teams/Zoom/Meet\nTranscript + wrap-up → PIR\nCost: ~6/30m + 2\nHuman join ≤15 min", fill=LIGHT_BLUE)

    card(s, Inches(0.4), Inches(4.5), Inches(3.8), Inches(1.9), "Shift Agent", "Level-1 OOO conflict\nSlack DM coverage OR\nPath B: Calendar + web override\nCost: 0 Actions", fill=LIGHT_GREEN, title_color=GREEN, line=GREEN)
    card(s, Inches(8.9), Inches(4.5), Inches(4.0), Inches(1.9), "Insights Agent", "MTTR / MTTA / volume Q&A\nSanity-check Analytics UI\nWeekly tips: Slack-first\nCost: 0 Actions", fill=LIGHT_GREEN, title_color=GREEN, line=GREEN)

    card(s, Inches(4.9), Inches(4.5), Inches(3.5), Inches(1.9), "Human", "Owns blast radius\nApprove remediations\nResolve → SRE memory", fill=WHITE, line=NAVY)
    footer(s, 7)

    # ----- 8 Checklist -----
    s = blank(prs)
    header(s, "Design Checklist — Before / During POC")
    rows_top = Inches(1.05)
    items = [
        ("Detect", "Dynatrace → PD test Integration Key only; problem URL in payload"),
        ("Noise (optional)", "Event Orchestration / grouping if alert storms exist (see seq 7)"),
        ("Advance", "AI Settings → Teams/Slack Connected → SRE/Scribe/Insights(+Shift) Enabled"),
        ("Identity", "Each POC user: linkUser; graphAuth if Delegated-heavy"),
        ("Safe target", "poc-pd-ai-agents-test Service; EP pages only you; poc-shift-agent Level-1"),
        ("SRE", "Runbook ≤100 KB; confirm remediations; Resolve to save memory"),
        ("Scribe", "Meeting ≤10 min POC; warn attendees; lobby admit"),
        ("Shift", "Teams-only → Path B web+Calendar; do not claim Slack DM worked"),
        ("Insights", "Ask in Teams; document weekly DM Slack gap"),
        ("Board pack", "Pair this diagram PPT with seq 14 architecture md + seq 13 narrative PPT"),
    ]
    y = rows_top
    for idx, (title, body) in enumerate(items):
        fill = WHITE if idx % 2 == 0 else LIGHT_GRAY
        card(s, Inches(0.35), y, Inches(12.6), Inches(0.52), title, body, fill=fill, line=NAVY_MID)
        y += Inches(0.55)
    footer(s, 8)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote {out}")


if __name__ == "__main__":
    base = Path(__file__).resolve().parent
    build(base / "15-Dynatrace-PagerDuty-Four-Agents-Design-Diagrams.pptx")
