#!/usr/bin/env python3
"""Combine seq 13 + 15 + 17 into one whole PPT with NO bottom footer / page numbers."""

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x2C, 0x5C)
LIGHT_BLUE = RGBColor(0xD6, 0xE6, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x22)
GREEN = RGBColor(0x1E, 0x84, 0x4E)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
GRAY = RGBColor(0x55, 0x55, 0x55)

SRC_ROOT = Path("/Users/k/Learnings/AIProject/CursorFiles/Daily Files/2026-08-31")
SRC_13 = SRC_ROOT / "13-dynatrace-pagerduty-four-agents-ppt/13-Dynatrace-PagerDuty-Four-AI-Agents-ARB-style.pptx"
SRC_15 = SRC_ROOT / "15-dynatrace-pd-four-agents-design-diagrams/15-Dynatrace-PagerDuty-Four-Agents-Design-Diagrams.pptx"
SRC_17 = SRC_ROOT / "17-pagerduty-four-agents-teams-only-poc-ppt/17-PagerDuty-Four-Agents-Teams-Only-POC.pptx"

SKIP_13 = {0}
SKIP_15 = {0}
SKIP_17 = {0}

FOOTER_MARKERS = (
    "AXA Japan ARB Template style",
    "ARB-style",
    "ARB style",
    "OneAXA",
    "Japan ARB",
    "LAST UPDATED",
    "PagerDuty Four Agents — Teams-Only POC",
    "Dynatrace to PagerDuty Four AI Agents — Design",
    "Dynatrace to PagerDuty Four AI Agents — Whole Pack",
    "Dynatrace → PagerDuty Four AI Agents",
    "ARB-style working deck",
    "Whole Pack",
)


def set_run(run, text, size=12, bold=False, color=BLACK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_textbox(slide, left, top, width, height, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color)
    return box


def add_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def copy_slide(dest_prs, source_slide):
    dest_slide = dest_prs.slides.add_slide(dest_prs.slide_layouts[6])
    sp_tree = dest_slide.shapes._spTree
    for child in list(sp_tree):
        tag = child.tag
        if tag.endswith("}sp") or tag.endswith("}pic") or tag.endswith("}graphicFrame") or tag.endswith("}cxnSp") or tag.endswith("}grpSp"):
            sp_tree.remove(child)
    for shape in source_slide.shapes:
        sp_tree.insert_element_before(deepcopy(shape.element), "p:extLst")
    if source_slide.has_notes_slide:
        notes_text = source_slide.notes_slide.notes_text_frame.text
        if notes_text.strip():
            dest_slide.notes_slide.notes_text_frame.text = notes_text
    return dest_slide


def strip_bottom_text(slide):
    """Remove every text box / short label near the bottom. No footer, no page numbers."""
    for shape in list(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = (shape.text_frame.text or "").strip()
        try:
            top = shape.top or 0
        except Exception:
            top = 0
        near_bottom = top > Inches(6.8)
        if not near_bottom:
            # Still scrub banned phrases if they appear anywhere as tiny leftover footers
            if any(m in text for m in FOOTER_MARKERS) and len(text) < 100:
                # only remove if it looks like a footer line, not a title
                try:
                    if top > Inches(6.5):
                        shape._element.getparent().remove(shape._element)
                except Exception:
                    pass
            continue

        # Bottom band: remove all short text (footers / page nums)
        if len(text) < 140 or text.isdigit() or " / " in text or any(m in text for m in FOOTER_MARKERS):
            shape._element.getparent().remove(shape._element)


def section_divider(prs, part, title, bullets, fill_accent=TEAL):
    s = blank(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_textbox(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.4), part, size=16, bold=True, color=fill_accent)
    add_textbox(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.8), title, size=32, bold=True, color=WHITE)
    y = Inches(3.5)
    for b in bullets:
        add_textbox(s, Inches(0.7), y, Inches(11.5), Inches(0.4), f"• {b}", size=16, color=LIGHT_BLUE)
        y += Inches(0.45)
    return s


def master_title(prs):
    s = blank(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.4), "Combined architecture + design + Teams POC pack", size=14, bold=True, color=LIGHT_BLUE)
    add_textbox(s, Inches(0.5), Inches(1.7), Inches(12), Inches(0.7), "Dynatrace to PagerDuty", size=36, bold=True, color=WHITE)
    add_textbox(s, Inches(0.5), Inches(2.5), Inches(12), Inches(0.5), "Four AI Agents — Whole Deck", size=26, bold=True, color=LIGHT_BLUE)
    add_textbox(s, Inches(0.5), Inches(3.2), Inches(12), Inches(0.4), "SRE  |  Scribe  |  Shift  |  Insights", size=16, color=WHITE)
    card_fill = RGBColor(0x12, 0x3A, 0x6B)
    for i, (title, body) in enumerate(
        [
            ("Part A — Narrative", "Architecture summary\nAS-IS / TO-BE story\nAgent detail + cost/safety"),
            ("Part B — Design diagrams", "Solution context boxes\nImpact / Alignment / HLD\nAgent hub + checklist"),
            ("Part C — Teams-only POC", "§0 enablement\nSRE / Scribe / Shift / Insights\nClick-path runbook"),
        ]
    ):
        left = Inches(0.5 + i * 4.2)
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(4.2), Inches(4.0), Inches(2.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = card_fill
        shape.line.color.rgb = LIGHT_BLUE
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_top = Inches(0.12)
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        set_run(r0, title, size=14, bold=True, color=WHITE)
        for line in body.split("\n"):
            p = tf.add_paragraph()
            r = p.add_run()
            set_run(r, line, size=12, color=LIGHT_BLUE)
    return s


def agenda_slide(prs, counts):
    s = blank(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.75), NAVY)
    add_textbox(s, Inches(0.35), Inches(0.18), Inches(12.5), Inches(0.45), "Whole Pack — Agenda", size=20, bold=True, color=WHITE)
    rows = [
        ("Part", "Source", "What you get", "Slides (approx)"),
        ("A", "seq 13 narrative PPT", "Summary sheet, AS-IS/TO-BE story, agents, enablement", str(counts["a"])),
        ("B", "seq 15 design diagrams PPT", "Visual Solution Context, Impact, Alignment, HLD, hub", str(counts["b"])),
        ("C", "seq 17 Teams-only POC PPT", "Full §0–§4 click path for Teams-only orgs", str(counts["c"])),
    ]
    table_shape = s.shapes.add_table(len(rows), 4, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.45 * len(rows)))
    table = table_shape.table
    widths = [Inches(1.2), Inches(3.8), Inches(5.5), Inches(2.0)]
    for i, w in enumerate(widths):
        table.columns[i].width = w
    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            is_h = r_idx == 0
            set_run(run, cell_text, size=12, bold=is_h, color=WHITE if is_h else BLACK)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if is_h else (RGBColor(0xF2, 0xF4, 0xF7) if r_idx % 2 else WHITE)
    add_textbox(
        s,
        Inches(0.4),
        Inches(4.0),
        Inches(12.5),
        Inches(2.0),
        "How to use\n"
        "• Design review / board discussion → Part A then Part B\n"
        "• Hands-on Teams-only POC afternoon → Part C\n"
        "• Safe demo rule everywhere: test Service + Level-1 test schedule — never prod pages\n"
        "• No footer text and no page numbers on any slide",
        size=14,
        color=BLACK,
    )
    return s


def append_from(src_path, dest_prs, skip_indices):
    src = Presentation(str(src_path))
    for i, slide in enumerate(src.slides):
        if i in skip_indices:
            continue
        copy_slide(dest_prs, slide)


def build(out: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    n13 = len(Presentation(str(SRC_13)).slides) - len(SKIP_13)
    n15 = len(Presentation(str(SRC_15)).slides) - len(SKIP_15)
    n17 = len(Presentation(str(SRC_17)).slides) - len(SKIP_17)

    master_title(prs)
    agenda_slide(prs, {"a": n13, "b": n15, "c": n17})

    section_divider(
        prs,
        "PART A",
        "Narrative Architecture",
        [
            "From seq 13 — Dynatrace → PagerDuty Four AI Agents",
            "Summary sheet, AS-IS / TO-BE, data flow, agent detail",
            "Enablement, cost, safety, POC afternoon",
        ],
        fill_accent=LIGHT_BLUE,
    )
    append_from(SRC_13, prs, SKIP_13)

    section_divider(
        prs,
        "PART B",
        "Design Diagrams",
        [
            "From seq 15 — Solution Context design pack",
            "AS-IS / TO-BE hub diagrams, Impacted Platforms",
            "Alignment stack, Technical HLD, Agent hub, checklist",
        ],
        fill_accent=GREEN,
    )
    append_from(SRC_15, prs, SKIP_15)

    section_divider(
        prs,
        "PART C",
        "Teams-Only POC Runbook",
        [
            "From seq 17 — detailed Teams-only click path",
            "§0 shared enablement, then SRE → Scribe → Shift → Insights",
            "Honest Slack-first gaps documented for Shift DMs and Insights weekly DMs",
        ],
        fill_accent=ORANGE,
    )
    append_from(SRC_17, prs, SKIP_17)

    for slide in prs.slides:
        strip_bottom_text(slide)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote {out}")
    print(f"Total slides: {len(prs.slides)} (no footers)")


if __name__ == "__main__":
    base = Path(__file__).resolve().parent
    build(base / "19-Dynatrace-PagerDuty-Four-Agents-Whole-Pack-NoFooter.pptx")
